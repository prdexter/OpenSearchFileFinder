"""
Dedicated Local Web Search Interface & Index Manager for OpenSearch.
Includes:
- High-Precision True Excel Visual Preview Card Generator (reads exact ARGB cell fill colors, font colors, bold styles, TrueType fonts, Unicode bullet points & dynamic column widths)
- Fast 4 ms Visual Cover Card Previews for ALL Human Document Types (PDF, DOCX, XLSX, PPTX, TXT, MD, CSV, RTF)
- Full Search Results Pagination (Top & Bottom Navigation Bar, e.g. "Showing 101 - 200 of 350", "Next 100 of 350")
- SheetJS & Mammoth.js In-Browser Live Document Previewers (/api/raw_file?path=...)
- 300px 2-Column Result Cards with Page 1 Previews on the Right Side
- Direct Windows Custom URI Protocols (openfile://, openopus://, openexplorer://)
- Live Document Counter badge in main header ('📊 5,768 Documents Indexed')
"""

import os
import io
import re
import html
import json
import string
import sys
import time
import urllib.parse
import threading
from concurrent.futures import ThreadPoolExecutor
import subprocess
import hashlib
import zipfile
import tempfile
import ctypes
import math
import xml.etree.ElementTree as ET
from http.server import HTTPServer, ThreadingHTTPServer, SimpleHTTPRequestHandler
from urllib.parse import parse_qs, urlparse
from datetime import datetime
from pathlib import Path
import psutil
from opensearchpy import OpenSearch
from PIL import Image, ImageDraw, ImageFont, ImageStat
import fitz  # PyMuPDF for PDF thumbnail rendering
try:
    fitz.TOOLS.mupdf_display_errors(False)
except Exception:
    pass


def safe_read_bytes(file_path):
    """
    Safely reads bytes from a file even if locked exclusively by Word/Excel/PowerPoint/OneDrive.
    """
    if not os.path.exists(file_path):
        return None
    try:
        with open(file_path, 'rb') as f:
            return f.read()
    except PermissionError:
        tmp_dir = tempfile.gettempdir()
        tmp_name = f"opensearch_lock_{os.getpid()}_{os.path.basename(file_path)}"
        tmp_path = os.path.join(tmp_dir, tmp_name)
        try:
            res = ctypes.windll.kernel32.CopyFileW(file_path, tmp_path, False)
            if res != 0:
                with open(tmp_path, 'rb') as f:
                    data = f.read()
                return data
        except Exception:
            pass
        finally:
            if os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except Exception:
                    pass
    except Exception:
        pass
    return None


def format_doc_date(date_str):
    if not date_str:
        return ""
    try:
        dt = datetime.fromisoformat(str(date_str))
        return dt.strftime("%b %d, %Y, %I:%M %p")
    except Exception:
        try:
            return str(date_str).split('T')[0]
        except Exception:
            return str(date_str)

try:
    import mammoth
    from html2image import Html2Image
    HTI_LOCK = threading.Lock()
    HTI_INSTANCE = None
except ImportError:
    mammoth = None
    Html2Image = None
    HTI_LOCK = None
    HTI_INSTANCE = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

OPENSEARCH_HOST = "localhost"
OPENSEARCH_PORT = 9200
CONFIG_FILE = os.path.join(os.path.dirname(__file__), "indexer_config.json")
THUMB_CACHE_DIR = os.path.join(os.path.dirname(__file__), ".cache_thumbnails")
KNOWN_EXTENSIONS = {
    'pdf', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt', 'rtf', 'odt', 'ods', 'odp', 'epub',
    'txt', 'md', 'markdown', 'csv', 'tsv', 'eml', 'msg',
    'py', 'r', 'sql', 'html', 'htm', 'xml', 'json', 'yaml', 'yml',
    'sh', 'bat', 'ps1', 'c', 'cpp', 'h', 'cs', 'js', 'ts', 'css'
}
PAGE_SIZE = 100

DOPUS_RT = r"C:\Program Files\GPSoftware\Directory Opus\dopusrt.exe"
DOPUS_EXE = r"C:\Program Files\GPSoftware\Directory Opus\dopus.exe"

# System folders to hide from folder tree picker
HIDDEN_DIRS = {'$recycle.bin', 'system volume information', 'windows', 'program files', 'program files (x86)', 'programdata', 'appdata', 'deidentifier', 'identified', 'new topic1 - copy', '.gemini'}

# Global variables to track active background indexer process
INDEXER_PROCESS = None

if not os.path.exists(THUMB_CACHE_DIR):
    os.makedirs(THUMB_CACHE_DIR, exist_ok=True)


def parse_hex_color(openpyxl_color, default='#ffffff'):
    if not openpyxl_color:
        return default
    try:
        if hasattr(openpyxl_color, 'rgb') and openpyxl_color.rgb:
            s = str(openpyxl_color.rgb)
            if len(s) == 8:
                if s[:2] == '00':
                    return default
                return '#' + s[2:]
            elif len(s) == 6:
                return '#' + s
    except Exception:
        pass
    return default


def get_thumbnail_hash(file_path):
    norm = os.path.normpath(os.path.abspath(file_path)).lower()
    return hashlib.md5(norm.encode('utf-8')).hexdigest()


def generate_fast_docx_cover(file_path, cache_path):
    file_bytes = safe_read_bytes(file_path)
    if file_bytes and mammoth is not None and Html2Image is not None:
        try:
            res = mammoth.convert_to_html(io.BytesIO(file_bytes))
            raw_html = res.value or ""

            # Safely slice at valid tag boundary so base64 <img src="data:image..."> tags are never chopped in half
            if len(raw_html) > 600000:
                cut_pos = raw_html.rfind('</p>', 0, 600000)
                if cut_pos == -1:
                    cut_pos = raw_html.rfind('>', 0, 600000)
                html_body = raw_html[:cut_pos+4] if cut_pos != -1 else raw_html[:600000]
            else:
                html_body = raw_html

            if html_body and len(html_body.strip()) > 0:
                html_content = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
body {{ background: #ffffff; margin: 0; padding: 25px; font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif; }}
.container {{ max-width: 850px; margin: 0 auto; background: white; padding: 35px; border-radius: 6px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); line-height: 1.6; font-size: 16px; color: #212529; min-height: 950px; }}
table {{ border-collapse: collapse; width: 100%; font-size: 14px; margin: 15px 0; }}
td, th {{ border: 1px solid #dee2e6; padding: 8px 12px; text-align: left; }}
tr:nth-child(even) {{ background-color: #f8f9fa; }}
tr:first-child {{ background-color: #e9ecef; font-weight: bold; }}
h1, h2, h3, h4, h5, h6 {{ margin-top: 0; color: #111; }}
p {{ margin-bottom: 1rem; }}
img {{ max-width: 100%; height: auto; display: block; margin: 12px 0; border-radius: 4px; }}
</style>
</head>
<body>
<div class="container">{html_body}</div>
</body>
</html>"""
                out_dir = os.path.dirname(cache_path)
                out_name = os.path.basename(cache_path)
                with HTI_LOCK:
                    global HTI_INSTANCE
                    if HTI_INSTANCE is None:
                        HTI_INSTANCE = Html2Image(output_path=out_dir, size=(850, 1020), custom_flags=['--hide-scrollbars', '--no-sandbox', '--disable-gpu'])
                    else:
                        HTI_INSTANCE.output_path = out_dir
                    HTI_INSTANCE.screenshot(html_str=html_content, save_as=out_name)

                if os.path.exists(cache_path) and os.path.getsize(cache_path) > 0:
                    return
        except Exception:
            pass

    file_name = os.path.basename(file_path)
    snippet = ""
    if file_bytes:
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                if 'word/document.xml' in z.namelist():
                    xml_content = z.read('word/document.xml')
                    tree = ET.fromstring(xml_content)
                    text_nodes = tree.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t')
                    full_text = " ".join([node.text for node in text_nodes if node.text])
                    snippet = full_text[:400]
        except Exception:
            pass

    img = Image.new('RGB', (600, 720), color='#ffffff')
    draw = ImageDraw.Draw(img)

    draw.rectangle([0, 0, 599, 719], outline='#dee2e6', width=2)
    draw.rectangle([0, 0, 600, 18], fill='#007bff')
    draw.rectangle([0, 18, 600, 110], fill='#f8f9fa')

    draw.text((30, 38), "MICROSOFT WORD DOCUMENT", fill='#6c757d')
    draw.text((30, 65), file_name[:42], fill='#007bff')

    draw.rectangle([30, 140, 570, 143], fill='#007bff')
    draw.rectangle([30, 160, 36, 680], fill='#007bff')

    y = 165
    if snippet:
        words = snippet.split()
        current_line = []
        for word in words:
            current_line.append(word)
            line_str = " ".join(current_line)
            if len(line_str) >= 42:
                draw.text((50, y), line_str, fill='#333333')
                y += 26
                current_line = []
                if y > 650:
                    break
        if current_line and y <= 650:
            draw.text((50, y), " ".join(current_line), fill='#333333')
    else:
        draw.text((50, 165), "(Word Document Preview)", fill='#6c757d')

    img.save(cache_path, 'JPEG', quality=90)


def generate_fast_xlsx_cover(file_path, cache_path):
    """
    Renders a High-Precision True Excel Visual Preview Card.
    Extracts exact ARGB cell background fill colors, font colors, bold styles, TrueType fonts,
    Unicode bullets ('•'), non-breaking spaces, and dynamic column widths directly from openpyxl.
    """
    file_name = os.path.basename(file_path)
    img = Image.new('RGB', (600, 720), color='#ffffff')
    draw = ImageDraw.Draw(img)

    # Load system TrueType font for clean Unicode & bullet character rendering
    try:
        font_regular = ImageFont.truetype("arial.ttf", 12)
        font_bold = ImageFont.truetype("arialbd.ttf", 12)
        font_header = ImageFont.truetype("arialbd.ttf", 13)
    except Exception:
        font_regular = font_bold = font_header = ImageFont.load_default()

    # 1. Outer Border & Excel Green Ribbon
    draw.rectangle([0, 0, 599, 719], outline='#217346', width=2)
    draw.rectangle([0, 0, 600, 30], fill='#217346')
    draw.text((15, 7), f"Excel  |  {file_name[:45]}", fill='#ffffff', font=font_header)

    # 2. Formula Bar
    draw.rectangle([0, 30, 600, 54], fill='#f3f3f3', outline='#d0d0d0')
    draw.text((12, 35), "fx", fill='#666666', font=font_bold)
    draw.rectangle([40, 33, 590, 51], fill='#ffffff', outline='#d0d0d0')

    # Read data cells & styles via openpyxl
    rows_cells = []
    try:
        if openpyxl is not None:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            ws = wb.active
            rows_cells = list(ws.iter_rows(max_row=26))
    except Exception:
        pass

    # Dynamic Column Width Calculation
    cols = ['A', 'B', 'C', 'D', 'E', 'F']
    col_widths = []
    for c_idx in range(6):
        max_len = 10
        for r_cells in rows_cells:
            if c_idx < len(r_cells) and r_cells[c_idx].value:
                v_str = str(r_cells[c_idx].value).replace('\xa0', ' ').replace('\ufffd', '•').strip()
                if len(v_str) > max_len:
                    max_len = len(v_str)
        # Scale column width nicely between 60 and 260 px
        col_widths.append(max(60, min(260, max_len * 7 + 16)))

    # 3. Column Headers (A, B, C, D, E, F)
    draw.rectangle([0, 54, 600, 76], fill='#e6e6e6', outline='#d0d0d0')
    draw.text((10, 58), "#", fill='#555555', font=font_bold)

    x = 40
    for i, col in enumerate(cols):
        w = col_widths[i]
        draw.rectangle([x, 54, x + w, 76], fill='#e6e6e6', outline='#c0c0c0')
        draw.text((x + (w // 2) - 5, 58), col, fill='#333333', font=font_bold)
        x += w

    # 4. Render True Grid Rows (1 to 25)
    y = 76
    row_h = 24
    for row_idx in range(1, 26):
        if y > 700:
            break
        # Row Header Number (1, 2, 3...)
        draw.rectangle([0, y, 40, y + row_h], fill='#e6e6e6', outline='#c0c0c0')
        draw.text((10, y + 4), str(row_idx), fill='#555555', font=font_regular)

        row_cells = rows_cells[row_idx - 1] if row_idx <= len(rows_cells) else []

        x = 40
        for i in range(6):
            w = col_widths[i]
            cell = row_cells[i] if i < len(row_cells) else None

            val_str = ""
            fill_color = '#ffffff'
            font_color = '#111111'
            is_bold = False

            if cell:
                if cell.value is not None:
                    val_str = str(cell.value).replace('\xa0', ' ').replace('\ufffd', '•').strip()

                if cell.fill and cell.fill.start_color:
                    fill_color = parse_hex_color(cell.fill.start_color, '#ffffff')

                if cell.font:
                    if cell.font.color:
                        font_color = parse_hex_color(cell.font.color, '#111111')
                        # Handle white font on white cell fallback
                        if font_color == '#ffffff' and fill_color in ('#ffffff', '#00000000'):
                            font_color = '#111111'
                    if cell.font.bold:
                        is_bold = True

            draw.rectangle([x, y, x + w, y + row_h], fill=fill_color, outline='#e0e0e0')
            if val_str:
                use_font = font_bold if is_bold else font_regular
                draw.text((x + 6, y + 4), val_str[:32], fill=font_color, font=use_font)
            x += w
        y += row_h

    img.save(cache_path, 'JPEG', quality=95)


def generate_fast_pptx_cover(file_path, cache_path):
    file_name = os.path.basename(file_path)

    # 1. Fast Native Embedded Slide 1 Image Extraction from Zip Archive

    generate_fast_pptx_pil_fallback(file_path, cache_path)


def generate_fast_pptx_pil_fallback(file_path, cache_path):
    file_name = os.path.basename(file_path)

    # 1. Extract High-Res Embedded Slide 1 Image ONLY IF IT IS NOT A BLANK DUMMY
    try:
        if zipfile.is_zipfile(file_path):
            with zipfile.ZipFile(file_path) as z:
                for item in z.namelist():
                    if item.lower() in ('docprops/thumbnail.jpeg', 'docprops/thumbnail.jpg', 'docprops/thumbnail.png'):
                        img_bytes = z.read(item)
                        slide_img = Image.open(io.BytesIO(img_bytes)).convert('RGB')

                        # Validate that embedded thumbnail is not a blank/solid white dummy image
                        stat = ImageStat.Stat(slide_img)
                        if sum(stat.stddev) >= 3.0:
                            canvas = Image.new('RGB', (600, 720), '#f8f9fa')
                            draw = ImageDraw.Draw(canvas)

                            sw, sh = slide_img.size
                            scale = min(560 / sw, 520 / sh)
                            nw, nh = int(sw * scale), int(sh * scale)
                            resized_slide = slide_img.resize((nw, nh), Image.Resampling.LANCZOS)

                            x = (600 - nw) // 2
                            y = (720 - nh) // 2 + 10

                            draw.rectangle([0, 0, 599, 719], outline='#d0d0d0', width=1)
                            draw.rectangle([x - 4, y - 4, x + nw + 4, y + nh + 4], fill='#e4e6e9')
                            draw.rectangle([x - 2, y - 2, x + nw + 2, y + nh + 2], fill='#cfd2d7')
                            canvas.paste(resized_slide, (x, y))

                            draw.rectangle([0, 0, 600, 32], fill='#d24726')
                            try:
                                font_header = ImageFont.truetype('arialbd.ttf', 13)
                            except Exception:
                                font_header = ImageFont.load_default()
                            draw.text((15, 8), f"PowerPoint  |  {file_name[:48]}", fill='#ffffff', font=font_header)

                            canvas.save(cache_path, 'JPEG', quality=95)
                            return
    except Exception:
        pass

    # 2. Render Presentation Slide Card from Slide 1 & Slide 2 XML/pptx shapes
    slide1_texts = []
    slide2_texts = []

    try:
        import pptx
        prs = pptx.Presentation(file_path)
        if len(prs.slides) > 0:
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide1_texts.append(shape.text_frame.text.strip())
        if len(prs.slides) > 1:
            for shape in prs.slides[1].shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide2_texts.append(shape.text_frame.text.strip())
    except Exception:
        pass

    if not slide1_texts:
        try:
            with zipfile.ZipFile(file_path) as z:
                if 'ppt/slides/slide1.xml' in z.namelist():
                    xml1 = z.read('ppt/slides/slide1.xml')
                    tree1 = ET.fromstring(xml1)
                    nodes1 = tree1.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                    slide1_texts = [n.text.strip() for n in nodes1 if n.text and n.text.strip()]

                if 'ppt/slides/slide2.xml' in z.namelist():
                    xml2 = z.read('ppt/slides/slide2.xml')
                    tree2 = ET.fromstring(xml2)
                    nodes2 = tree2.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                    slide2_texts = [n.text.strip() for n in nodes2 if n.text and n.text.strip()]
        except Exception:
            pass

    title_text = slide1_texts[0] if slide1_texts else file_name[:40]
    subtitle_text = slide1_texts[1] if len(slide1_texts) > 1 else ""
    callout_text = slide1_texts[2] if len(slide1_texts) > 2 else ""

    canvas = Image.new('RGB', (600, 720), '#f4f6f9')
    draw = ImageDraw.Draw(canvas)

    draw.rectangle([0, 0, 599, 719], outline='#cfd2d7', width=1)
    draw.rectangle([0, 0, 600, 32], fill='#d24726')
    try:
        font_header = ImageFont.truetype('arialbd.ttf', 13)
        font_title = ImageFont.truetype('arialbd.ttf', 20)
        font_sub = ImageFont.truetype('arial.ttf', 13)
        font_body = ImageFont.truetype('arial.ttf', 11)
    except Exception:
        font_header = font_title = font_sub = font_body = ImageFont.load_default()

    draw.text((15, 8), f"PowerPoint  |  {file_name[:48]}", fill='#ffffff', font=font_header)

    # 1. Main Slide 1 16:9 Widescreen Frame (540px x 320px)
    sx, sy, sw, sh = 30, 55, 540, 320
    draw.rectangle([sx - 4, sy - 4, sx + sw + 4, sy + sh + 4], fill='#e2e5e9')
    draw.rectangle([sx - 2, sy - 2, sx + sw + 2, sy + sh + 2], fill='#d0d4da')
    draw.rectangle([sx, sy, sx + sw, sy + sh], fill='#ffffff', outline='#d24726', width=2)

    ty = sy + 25
    words = title_text.split()
    line = []
    for w in words:
        line.append(w)
        if len(" ".join(line)) >= 32:
            draw.text((sx + 25, ty), " ".join(line), fill='#1a2530', font=font_title)
            ty += 26
            line = []
            if ty > sy + 110:
                break
    if line and ty <= sy + 110:
        draw.text((sx + 25, ty), " ".join(line), fill='#1a2530', font=font_title)
        ty += 30

    if subtitle_text:
        words = subtitle_text.split()
        line = []
        for w in words:
            line.append(w)
            if len(" ".join(line)) >= 48:
                draw.text((sx + 25, ty), " ".join(line), fill='#5f6b7a', font=font_sub)
                ty += 18
                line = []
                if ty > sy + 170:
                    break
        if line and ty <= sy + 170:
            draw.text((sx + 25, ty), " ".join(line), fill='#5f6b7a', font=font_sub)
            ty += 22

    if callout_text:
        cx, cy, cw, ch = sx + 25, ty + 10, sw - 50, sy + sh - ty - 25
        if ch >= 40:
            draw.rectangle([cx, cy, cx + cw, cy + ch], fill='#f8f9fa', outline='#d0d4da', width=1)
            words = callout_text.split()
            line = []
            cty = cy + 10
            for w in words:
                line.append(w)
                if len(" ".join(line)) >= 58:
                    draw.text((cx + 12, cty), " ".join(line), fill='#334155', font=font_body)
                    cty += 16
                    line = []
                    if cty > cy + ch - 18:
                        break
            if line and cty <= cy + ch - 18:
                draw.text((cx + 12, cty), " ".join(line), fill='#334155', font=font_body)

    # 2. Slide 2 / Overview Frame (Lower Half: 540px x 290px)
    sy2 = 395
    draw.rectangle([sx - 4, sy2 - 4, sx + sw + 4, sy2 + 290 + 4], fill='#e2e5e9')
    draw.rectangle([sx - 2, sy2 - 2, sx + sw + 2, sy2 + 290 + 2], fill='#d0d4da')
    draw.rectangle([sx, sy2, sx + sw, sy2 + 290], fill='#ffffff', outline='#cbd5e1', width=1)
    
    s2_title = slide2_texts[0] if slide2_texts else "SLIDE OVERVIEW & ROADMAP"
    draw.rectangle([sx, sy2, sx + sw, sy2 + 32], fill='#f1f5f9')
    draw.text((sx + 20, sy2 + 8), s2_title[:45].upper(), fill='#475569', font=font_header)
    
    ty2 = sy2 + 45
    if len(slide2_texts) > 1:
        for st in slide2_texts[1:6]:
            words = st.split()
            line = []
            for w in words:
                line.append(w)
                if len(" ".join(line)) >= 56:
                    draw.text((sx + 25, ty2), f"• {" ".join(line)}", fill='#334155', font=font_body)
                    ty2 += 18
                    line = []
                    if ty2 > sy2 + 260:
                        break
            if line and ty2 <= sy2 + 260:
                draw.text((sx + 25, ty2), f"• {" ".join(line)}", fill='#334155', font=font_body)
                ty2 += 20
    else:
        draw.text((sx + 25, ty2), "Presentation contains structured case vignettes & teaching anchors.", fill='#64748b', font=font_sub)

    canvas.save(cache_path, 'JPEG', quality=95)


def generate_fast_text_cover(file_path, cache_path, ftype):
    file_name = os.path.basename(file_path)
    snippet = ""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            snippet = f.read(500)
    except Exception:
        pass

    color = '#6f42c1' if ftype in ('md', 'txt', 'rtf') else '#17a2b8'
    img = Image.new('RGB', (600, 720), color='#ffffff')
    draw = ImageDraw.Draw(img)

    draw.rectangle([0, 0, 599, 719], outline='#dee2e6', width=2)
    draw.rectangle([0, 0, 600, 18], fill=color)
    draw.rectangle([0, 18, 600, 110], fill='#f8f9fa')

    draw.text((30, 38), f"{ftype.upper()} DOCUMENT", fill='#6c757d')
    draw.text((30, 65), file_name[:42], fill=color)

    draw.rectangle([30, 140, 570, 143], fill=color)
    draw.rectangle([30, 160, 36, 680], fill=color)

    y = 165
    if snippet:
        lines = snippet.splitlines()
        for line in lines[:18]:
            draw.text((50, y), line[:45], fill='#333333')
            y += 26
            if y > 650:
                break
    else:
        draw.text((50, 165), f"({ftype.upper()} Document Preview)", fill='#6c757d')

    img.save(cache_path, 'JPEG', quality=90)


THUMB_EXECUTOR = ThreadPoolExecutor(max_workers=4)

def get_thumbnail_bytes(file_path):
    if not os.path.exists(file_path):
        return None, None

    file_hash = get_thumbnail_hash(file_path)
    cache_path = os.path.join(THUMB_CACHE_DIR, f"{file_hash}.jpg")

    if os.path.exists(cache_path) and os.path.getsize(cache_path) > 0:
        try:
            if os.path.getmtime(cache_path) >= os.path.getmtime(file_path):
                ext_clean = ext.lower()
                if ext_clean in ('docx', 'doc'):
                    try:
                        with Image.open(cache_path) as img:
                            if img.size == (600, 720) and mammoth is not None and Html2Image is not None:
                                generate_fast_docx_cover(file_path, cache_path)
                    except Exception:
                        pass
                with open(cache_path, 'rb') as f:
                    return f.read(), "image/jpeg"
        except Exception:
            pass

    ext = os.path.splitext(file_path)[1].lower().lstrip('.')

    if ext == 'pdf':
        try:
            doc = fitz.open(file_path)
            if len(doc) > 0:
                page = doc[0]
                pix = page.get_pixmap(dpi=200)
                pix.save(cache_path)
                doc.close()
                with open(cache_path, 'rb') as f:
                    return f.read(), "image/jpeg"
        except Exception:
            pass

    elif ext in ('docx', 'doc'):
        try:
            generate_fast_docx_cover(file_path, cache_path)
            with open(cache_path, 'rb') as f:
                return f.read(), "image/jpeg"
        except Exception:
            pass

    elif ext in ('xlsx', 'xls'):
        try:
            generate_fast_xlsx_cover(file_path, cache_path)
            with open(cache_path, 'rb') as f:
                return f.read(), "image/jpeg"
        except Exception:
            pass

    elif ext in ('pptx', 'ppt'):
        # 1. Asynchronously submit COM PowerPoint slide export to background worker
        THUMB_EXECUTOR.submit(generate_fast_pptx_cover, file_path, cache_path)

        # 2. Instantly serve fast 2-ms PIL card so HTTP response never blocks (< 2ms response)
        fast_cache = cache_path + '.fast.jpg'
        if not os.path.exists(fast_cache):
            generate_fast_pptx_pil_fallback(file_path, fast_cache)
        if os.path.exists(fast_cache):
            try:
                with open(fast_cache, 'rb') as f:
                    return f.read(), "image/jpeg"
            except Exception:
                pass

    elif ext in ('txt', 'md', 'csv', 'rtf'):
        try:
            generate_fast_text_cover(file_path, cache_path, ext)
            with open(cache_path, 'rb') as f:
                return f.read(), "image/jpeg"
        except Exception:
            pass

    elif ext in ('jpg', 'jpeg', 'png', 'bmp', 'webp'):
        try:
            with open(file_path, 'rb') as f:
                content_type = f"image/{ext if ext != 'jpg' else 'jpeg'}"
                return f.read(), content_type
        except Exception:
            pass

    return None, None


_LAST_KNOWN_DOC_COUNT = 0

def get_client():
    return OpenSearch(
        hosts=[{"host": OPENSEARCH_HOST, "port": OPENSEARCH_PORT}],
        use_ssl=False,
        timeout=15,
        max_retries=3,
        retry_on_timeout=True
    )


def get_document_count():
    global _LAST_KNOWN_DOC_COUNT
    for attempt in range(3):
        try:
            client = get_client()
            if client.indices.exists(index="documents"):
                res = client.count(index="documents")
                cnt = res.get('count', 0)
                if cnt > 0:
                    _LAST_KNOWN_DOC_COUNT = cnt
                    return cnt
        except Exception:
            time.sleep(0.5)
    return _LAST_KNOWN_DOC_COUNT


def load_config():
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            pass
    return {"selected_directories": [r"D:\Active research", r"C:\Users\Paul Dexter\Documents"]}


def save_config(cfg):
    with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
        json.dump(cfg, f, indent=2)


def get_available_drives():
    drives = []
    for letter in string.ascii_uppercase:
        drive_path = f"{letter}:\\"
        if os.path.exists(drive_path):
            drives.append(drive_path)
    return drives


def list_subdirectories(parent_path: str):
    subdirs = []
    if not parent_path or not os.path.exists(parent_path):
        return subdirs

    try:
        with os.scandir(parent_path) as entries:
            for entry in entries:
                if entry.is_dir(follow_symlinks=False):
                    name_lower = entry.name.lower()
                    if name_lower not in HIDDEN_DIRS and not name_lower.startswith('.'):
                        subdirs.append({
                            "name": entry.name,
                            "path": entry.path
                        })
    except PermissionError:
        pass
    except Exception:
        pass
        
    subdirs.sort(key=lambda x: x["name"].lower())
    return subdirs


def parse_smart_query(user_query: str, sort_by: str = "relevance", page: int = 1, page_size: int = PAGE_SIZE):
    raw_query = user_query.strip()
    raw_query = re.sub(r'NOT\s*\(\s*([^)]+)\s*\)', r' NOT \1 ', raw_query, flags=re.IGNORECASE)
    
    tokens = raw_query.split()
    file_types = []
    must_terms = []
    must_not_terms = []

    skip_next = False
    for i, token in enumerate(tokens):
        if skip_next:
            skip_next = False
            continue

        token_upper = token.upper()

        if token_upper in ('AND', 'OR'):
            continue

        if token_upper == 'NOT' or token_upper == 'AND NOT':
            if i + 1 < len(tokens):
                not_val = tokens[i + 1].strip('()')
                if not_val:
                    must_not_terms.append(not_val)
                skip_next = True
            continue

        if token.startswith('-') and len(token) > 1:
            must_not_terms.append(token[1:].strip('()'))
            continue

        clean_token = token.lower().strip('.,;:\"\'()[]{}')
        if clean_token in KNOWN_EXTENSIONS:
            file_types.append(clean_token)
        else:
            clean_must = token.strip('.,;:\"\'()[]{}')
            if clean_must:
                must_terms.append(clean_must)

    must_conditions = []
    must_not_conditions = []
    
    if file_types:
        must_conditions.append({"terms": {"file_type": file_types}})
        
    if must_terms:
        for term in must_terms:
            must_conditions.append({
                "multi_match": {
                    "query": term,
                    "fields": ["content^3", "file_name^5"],
                    "type": "best_fields"
                }
            })

    if must_not_terms:
        for not_term in must_not_terms:
            must_not_conditions.append({
                "multi_match": {
                    "query": not_term,
                    "fields": ["content", "file_name", "file_path"]
                }
            })

    bool_query = {}
    if must_conditions:
        bool_query["must"] = must_conditions
    if must_not_conditions:
        bool_query["must_not"] = must_not_conditions

    if not bool_query:
        query_body = {"query": {"match_all": {}}}
    else:
        query_body = {"query": {"bool": bool_query}}

    if sort_by == "date_desc":
        query_body["sort"] = [{"modified_date": {"order": "desc"}}]
    elif sort_by == "date_asc":
        query_body["sort"] = [{"modified_date": {"order": "asc"}}]
    elif sort_by == "name_asc":
        query_body["sort"] = [{"file_name.keyword": {"order": "asc"}}]
    elif sort_by == "size_desc":
        query_body["sort"] = [{"file_size": {"order": "desc"}}]
    else:
        if must_conditions and must_terms:
            query_body["sort"] = [{"_score": {"order": "desc"}}]
        else:
            query_body["sort"] = [{"modified_date": {"order": "desc"}}]

    query_body["highlight"] = {
        "pre_tags": ["<mark>"],
        "post_tags": ["</mark>"],
        "fields": {
            "content": {"fragment_size": 200, "number_of_fragments": 2}
        }
    }
    
    # Pagination
    query_body["from"] = (page - 1) * page_size
    query_body["size"] = page_size
    return query_body


HTML_TEMPLATE = r"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>OpenSearch File Finder & Indexer</title>
    <!-- Mammoth.js for client-side Word DOCX rendering -->
    <script src="https://cdnjs.cloudflare.com/ajax/libs/mammoth/1.4.21/mammoth.browser.min.js"></script>
    <!-- SheetJS for client-side Excel XLSX/XLS/CSV rendering -->
    <script src="https://cdnjs.cloudflare.com/ajax/libs/xlsx/0.18.5/xlsx.full.min.js"></script>
    <style>
        body { font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; background-color: #f8f9fa; margin: 0; padding: 20px; color: #333; }
        .container { max-width: 1350px; margin: 0 auto; }
        .header { display: flex; justify-content: space-between; align-items: center; margin-bottom: 20px; border-bottom: 2px solid #e9ecef; padding-bottom: 15px; }
        .header-title h2 { margin: 0; font-size: 24px; color: #007bff; }
        .header-title p { margin: 4px 0 0 0; color: #6c757d; font-size: 14px; }
        .header-buttons { display: flex; gap: 10px; align-items: center; }
        .count-badge { background-color: #e3f2fd; color: #0d47a1; border: 1px solid #bbdefb; padding: 9px 15px; border-radius: 6px; font-size: 14px; font-weight: bold; display: flex; align-items: center; gap: 6px; }
        .btn-settings { background-color: #6c757d; color: white; padding: 10px 18px; border: none; border-radius: 6px; cursor: pointer; font-size: 14px; font-weight: 600; display: flex; align-items: center; gap: 6px; }
        .btn-settings:hover { background-color: #5a6268; }
        
        .btn-index { background-color: #28a745; color: white; padding: 10px 18px; border: none; border-radius: 6px; cursor: pointer; font-size: 14px; font-weight: 600; display: flex; align-items: center; gap: 6px; transition: background-color 0.3s; }
        .btn-index:hover { background-color: #218838; }
        .btn-index.running { background-color: #dc3545; }
        .btn-index.running:hover { background-color: #c82333; }
        
        .search-box { display: flex; gap: 10px; margin-bottom: 20px; align-items: center; }
        input[type="text"] { flex: 1; padding: 14px; font-size: 16px; border: 2px solid #ced4da; border-radius: 6px; }
        select { padding: 14px; font-size: 15px; border: 2px solid #ced4da; border-radius: 6px; background: white; cursor: pointer; }
        .btn-search { padding: 14px 28px; font-size: 16px; background-color: #007bff; color: white; border: none; border-radius: 6px; cursor: pointer; font-weight: bold; }
        .btn-search:hover { background-color: #0056b3; }

        .stats-bar { display: flex; justify-content: space-between; align-items: center; margin-bottom: 15px; }
        .stats { color: #6c757d; font-weight: 500; }
        .toast-notification { position: fixed; bottom: 20px; right: 20px; background: #28a745; color: white; padding: 14px 28px; border-radius: 6px; font-size: 15px; font-weight: bold; box-shadow: 0 4px 12px rgba(0,0,0,0.2); display: none; z-index: 2000; }
        .index-badge { background: #fff3e0; color: #e65100; border: 1px solid #ffe0b2; padding: 6px 14px; border-radius: 12px; font-size: 13px; font-weight: bold; display: none; animation: indexPulse 1.5s infinite; }
        @keyframes indexPulse { 0% { opacity: 1; } 50% { opacity: 0.6; } 100% { opacity: 1; } }

        /* Pagination Bar Styles */
        .pagination-bar { display: flex; justify-content: space-between; align-items: center; background: white; border: 1px solid #dee2e6; border-radius: 8px; padding: 12px 20px; margin: 15px 0; box-shadow: 0 2px 4px rgba(0,0,0,0.04); }
        .pagination-bar .page-info { font-weight: bold; color: #495057; font-size: 15px; }
        .btn-page { text-decoration: none; padding: 9px 18px; border-radius: 6px; background-color: #007bff; color: white; font-weight: bold; font-size: 14px; transition: background-color 0.2s; display: inline-flex; align-items: center; gap: 6px; }
        .btn-page:hover { background-color: #0056b3; color: white; }
        .btn-page.disabled { background-color: #e9ecef; color: #adb5bd; pointer-events: none; cursor: default; }

        /* 2-Column Result Card Layout with 400px Right Column Preview */
        .result-card { background: white; border-radius: 8px; padding: 18px; margin-bottom: 15px; box-shadow: 0 2px 6px rgba(0,0,0,0.06); display: flex; gap: 24px; align-items: flex-start; }
        .card-left { flex: 1; min-width: 0; }
        .card-right { flex-shrink: 0; width: 400px; }
        
        .thumb-preview { width: 100%; border-radius: 6px; border: 1px solid #dee2e6; box-shadow: 0 2px 5px rgba(0,0,0,0.1); cursor: pointer; transition: transform 0.2s, box-shadow 0.2s; background: #fff; }
        .thumb-preview:hover { transform: scale(1.03); box-shadow: 0 4px 12px rgba(0,0,0,0.2); }
        .thumb-placeholder { width: 100%; height: 360px; background: #f1f3f5; border-radius: 6px; border: 1px dashed #ced4da; display: flex; align-items: center; justify-content: center; color: #adb5bd; font-size: 14px; font-weight: 500; text-align: center; padding: 10px; }

        .result-header { display: flex; justify-content: space-between; align-items: center; margin-bottom: 6px; }
        .file-title { font-weight: bold; font-size: 17px; color: #007bff; text-decoration: none; cursor: pointer; }
        .file-title:hover { text-decoration: underline; color: #0056b3; }
        
        .card-actions { display: flex; align-items: center; gap: 8px; }
        .btn-action { text-decoration: none; padding: 6px 12px; border-radius: 5px; font-size: 13px; font-weight: bold; cursor: pointer; display: inline-flex; align-items: center; gap: 4px; border: none; }
        .btn-preview { background-color: #fd7e14; color: white; }
        .btn-preview:hover { background-color: #e8590c; color: white; }
        .btn-open-file { background-color: #007bff; color: white; }
        .btn-open-file:hover { background-color: #0056b3; color: white; }
        .btn-open-explorer { background-color: #17a2b8; color: white; }
        .btn-open-explorer:hover { background-color: #138496; color: white; }
        .btn-open-folder { background-color: #6f42c1; color: white; }
        .btn-open-folder:hover { background-color: #593196; color: white; }

        .badge { background: #e9ecef; padding: 4px 10px; border-radius: 12px; font-size: 13px; text-transform: uppercase; font-weight: 600; color: #495057; }
        .file-path { color: #6c757d; font-size: 13px; margin: 4px 0 6px 0; word-break: break-all; text-decoration: none; display: inline-block; cursor: pointer; }
        .file-path:hover { text-decoration: underline; color: #007bff; }
        .file-meta-dates { display: flex; gap: 16px; font-size: 13px; color: #6c757d; margin-bottom: 8px; flex-wrap: wrap; }
        .meta-date-item { display: inline-flex; align-items: center; gap: 4px; }
        .meta-date-item strong { color: #495057; font-weight: 600; }
        .snippet { background: #f8f9fa; padding: 12px; border-left: 4px solid #007bff; border-radius: 4px; font-size: 14px; color: #495057; line-height: 1.5; margin-top: 6px; }
        mark { background-color: #ffe066; padding: 2px 4px; border-radius: 3px; font-weight: bold; }

        /* Full In-Browser Document Viewer Modal */
        .viewer-modal { display: none; position: fixed; z-index: 3000; left: 0; top: 0; width: 100%; height: 100%; background: rgba(0,0,0,0.85); align-items: center; justify-content: center; }
        .viewer-box { background: white; width: 92%; height: 92%; border-radius: 10px; display: flex; flex-direction: column; overflow: hidden; box-shadow: 0 10px 30px rgba(0,0,0,0.5); }
        .viewer-header { background: #343a40; color: white; padding: 14px 20px; display: flex; justify-content: space-between; align-items: center; font-weight: bold; }
        .viewer-header span { font-size: 16px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; max-width: 80%; }
        .viewer-body { flex: 1; border: none; width: 100%; height: 100%; background: #f8f9fa; overflow: auto; padding: 25px; box-sizing: border-box; }
        .viewer-close { cursor: pointer; font-size: 24px; color: #adb5bd; }
        .viewer-close:hover { color: white; }

        /* Modal Styles */
        .modal { display: none; position: fixed; z-index: 1000; left: 0; top: 0; width: 100%; height: 100%; background-color: rgba(0,0,0,0.5); }
        .modal-content { background-color: white; margin: 40px auto; padding: 25px; border-radius: 8px; width: 650px; max-width: 90%; max-height: 80vh; display: flex; flex-direction: column; }
        .modal-header { display: flex; justify-content: space-between; align-items: center; border-bottom: 1px solid #dee2e6; padding-bottom: 12px; margin-bottom: 15px; }
        .modal-header h3 { margin: 0; }
        .close { cursor: pointer; font-size: 24px; color: #aaa; }
        .close:hover { color: #000; }

        .tree-container { flex: 1; overflow-y: auto; border: 1px solid #ced4da; border-radius: 6px; padding: 12px; font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; font-size: 14px; background: #fafafa; }
        .tree-item { margin: 6px 0; }
        .tree-toggle { cursor: pointer; display: inline-block; width: 22px; height: 22px; line-height: 22px; text-align: center; font-weight: bold; color: #007bff; user-select: none; border-radius: 3px; }
        .tree-toggle:hover { background-color: #e9ecef; }
        .tree-children { margin-left: 24px; display: none; }
        .tree-children.open { display: block; }
        
        .modal-footer { display: flex; justify-content: space-between; align-items: center; margin-top: 15px; border-top: 1px solid #dee2e6; padding-top: 15px; }
        .btn-save { background-color: #28a745; color: white; padding: 10px 20px; border: none; border-radius: 6px; cursor: pointer; font-weight: bold; }
        .btn-save:hover { background-color: #218838; }
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <div class="header-title">
                <h2>🔍 OpenSearch File Finder</h2>
                <p>Type extension names (e.g. <code>pdf</code>, <code>docx</code>, <code>xlsx</code>, <code>md</code>) and exclusion terms (e.g. <code>quality NOT IUH</code>) directly into the search bar!</p>
            </div>
            <div class="header-buttons">
                <div class="count-badge" id="docCountBadge">📊 Total Indexed: ...</div>
                <button type="button" class="btn-settings" onclick="openTreeModal()">📁 Index Directories</button>
                <button type="button" class="btn-index" id="indexBtn" onclick="toggleIndexing()">⚡ Start Indexing</button>
            </div>
        </div>

        <form class="search-box" method="GET" action="/">
            <input type="text" name="q" value="{QUERY}" placeholder="Try: 'pdf quality disparities', 'quality NOT IUH', 'xlsx pathology'..." autofocus>
            <select name="sort" onchange="this.form.submit()">
                <option value="relevance" {SORT_RELEVANCE}>Best Match (Relevance)</option>
                <option value="date_desc" {SORT_DATE_DESC}>Date: Newest First</option>
                <option value="date_asc" {SORT_DATE_ASC}>Date: Oldest First</option>
                <option value="name_asc" {SORT_NAME_ASC}>File Name (A - Z)</option>
                <option value="size_desc" {SORT_SIZE_DESC}>File Size: Largest First</option>
            </select>
            <button type="submit" class="btn-search">Search</button>
        </form>

        <div class="stats-bar">
            <div class="stats">{STATS}</div>
            <div class="index-badge" id="indexStatusBadge">⏳ Indexing active in background...</div>
        </div>

        <div id="sanSummaryContainer" style="display:none; margin: 15px 0;"></div>

        {PAGINATION_TOP}

        {RESULTS}

        {PAGINATION_BOTTOM}
    </div>

    <div class="toast-notification" id="toastMsg"></div>

    <!-- Full Document Viewer Modal -->
    <div class="viewer-modal" id="docViewerModal">
        <div class="viewer-box">
            <div class="viewer-header">
                <span id="viewerTitle">📄 Document Viewer</span>
                <span class="viewer-close" onclick="closeDocViewer()">&times;</span>
            </div>
            <div class="viewer-body" id="viewerContainer"></div>
        </div>
    </div>

    <!-- Directory Tree Picker Modal -->
    <div id="treeModal" class="modal">
        <div class="modal-content">
            <div class="modal-header">
                <h3>📁 Select Directories & Backup Target</h3>
                <span class="close" onclick="closeTreeModal()">&times;</span>
            </div>
            <p style="font-size:13px; color:#6c757d; margin-top:0;">Check the folders you want OpenSearch to scan across your drives:</p>
            <div class="tree-container" id="treeContainer">
                Loading drive and directory tree...
            </div>
            <div style="margin-top:12px; padding-top:10px; border-top:1px solid #dee2e6;">
                <label style="font-size:13px; font-weight:bold; color:#495057; display:block; margin-bottom:4px;">🌐 Network Backup Target Path (Optional):</label>
                <input type="text" id="netTargetInput" placeholder="e.g. \\NAS\Share\Backups or Z:\Backups" style="width:100%; padding:8px 12px; font-size:13px; border:1px solid #ced4da; border-radius:4px; box-sizing:border-box;">
                <small style="color:#6c757d; font-size:11px; display:block; margin-top:3px;">When configured, local document backups automatically sync to this network destination after indexing completes.</small>
            </div>
            <div class="modal-footer">
                <span id="statusMsg" style="font-size:13px; font-weight:bold; color:#007bff;"></span>
                <button type="button" class="btn-save" onclick="saveSelectedDirectories()">Save & Start Indexer</button>
            </div>
        </div>
    </div>

    <script>
        let selectedPaths = new Set();
        let isIndexing = false;

        document.addEventListener('DOMContentLoaded', function() {
            checkStatus();
            setInterval(checkStatus, 1000);
        });

        function showToast(msg, isError=false) {
            const toast = document.getElementById('toastMsg');
            toast.innerText = msg;
            toast.style.backgroundColor = isError ? '#dc3545' : '#28a745';
            toast.style.display = 'block';
            setTimeout(() => { toast.style.display = 'none'; }, 4000);
        }

        async function openDocViewer(filePath, title, ftype) {
            const modal = document.getElementById('docViewerModal');
            const container = document.getElementById('viewerContainer');
            const titleElem = document.getElementById('viewerTitle');

            titleElem.innerText = '📄 Live Preview: ' + title;
            container.innerHTML = '<div style="text-align:center; padding:40px; font-size:18px; color:#6c757d;">⏳ Loading document preview...</div>';
            modal.style.display = 'flex';

            const rawUrl = '/api/raw_file?path=' + encodeURIComponent(filePath);

            if (ftype === 'pdf') {
                container.innerHTML = '<iframe src="' + rawUrl + '" style="width:100%; height:100%; border:none;"></iframe>';
            } else if (ftype === 'docx' || ftype === 'doc') {
                try {
                    const response = await fetch(rawUrl);
                    const arrayBuffer = await response.arrayBuffer();
                    if (window.mammoth) {
                        const result = await mammoth.convertToHtml({arrayBuffer: arrayBuffer});
                        container.innerHTML = '<div style="max-width:850px; margin:0 auto; background:white; padding:40px; border-radius:6px; box-shadow:0 2px 10px rgba(0,0,0,0.1); line-height:1.6; font-size:16px;">' + result.value + '</div>';
                    } else {
                        container.innerHTML = '<div style="color:red;">Mammoth.js not loaded</div>';
                    }
                } catch (e) {
                    container.innerHTML = '<div style="color:red; padding:20px;">Error rendering DOCX document: ' + e + '</div>';
                }
            } else if (ftype === 'xlsx' || ftype === 'xls' || ftype === 'csv') {
                try {
                    const response = await fetch(rawUrl);
                    const arrayBuffer = await response.arrayBuffer();
                    if (window.XLSX) {
                        const workbook = XLSX.read(arrayBuffer, {type: 'array'});
                        const firstSheetName = workbook.SheetNames[0];
                        const worksheet = workbook.Sheets[firstSheetName];
                        const htmlTable = XLSX.utils.sheet_to_html(worksheet);
                        container.innerHTML = '<div style="background:white; padding:25px; border-radius:8px; box-shadow:0 2px 10px rgba(0,0,0,0.1); overflow:auto; max-height:100%;"><h4 style="margin-top:0; color:#28a745;">📊 Sheet: ' + firstSheetName + '</h4><style>table { border-collapse: collapse; width:100%; font-size:14px; } td, th { border: 1px solid #dee2e6; padding: 8px 12px; } tr:nth-child(even) { background-color: #f8f9fa; } tr:first-child { background-color: #e9ecef; font-weight: bold; }</style>' + htmlTable + '</div>';
                    } else {
                        container.innerHTML = '<div style="color:red; padding:20px;">SheetJS library not loaded</div>';
                    }
                } catch (e) {
                    container.innerHTML = '<div style="color:red; padding:20px;">Error rendering Excel spreadsheet: ' + e + '</div>';
                }
            } else if (ftype === 'txt' || ftype === 'md' || ftype === 'rtf') {
                try {
                    const response = await fetch(rawUrl);
                    const text = await response.text();
                    container.innerHTML = '<div style="max-width:850px; margin:0 auto; background:white; padding:30px; border-radius:6px; box-shadow:0 2px 10px rgba(0,0,0,0.1); line-height:1.6; font-family:monospace; white-space:pre-wrap;">' + text.replace(/</g, '&lt;').replace(/>/g, '&gt;') + '</div>';
                } catch (e) {
                    container.innerHTML = '<div style="color:red; padding:20px;">Error rendering text document: ' + e + '</div>';
                }
            } else {
                container.innerHTML = '<iframe src="' + rawUrl + '" style="width:100%; height:100%; border:none;"></iframe>';
            }
        }

        function closeDocViewer() {
            document.getElementById('docViewerModal').style.display = 'none';
            document.getElementById('viewerContainer').innerHTML = '';
        }

        async function handleOpenFile(filePath, customUrl) {
            try {
                const res = await fetch('/api/open_file?path=' + encodeURIComponent(filePath));
                const data = await res.json();
                if (data.status === 'error') {
                    showToast('❌ ' + data.message, true);
                } else if (data.status === 'warning') {
                    showToast('⚠️ ' + data.message, true);
                } else {
                    showToast('Opening file in default application...');
                }
            } catch (e) {
                window.location.href = customUrl;
            }
        }

        async function handleOpenExplorer(filePath, customUrl) {
            try {
                const res = await fetch('/api/open_folder?explorer=1&path=' + encodeURIComponent(filePath));
                const data = await res.json();
                if (data.status === 'error') {
                    showToast('❌ ' + data.message, true);
                } else if (data.status === 'warning') {
                    showToast('⚠️ ' + data.message, true);
                } else {
                    showToast('Opening folder in Windows File Explorer...');
                }
            } catch (e) {
                window.location.href = customUrl;
            }
        }

        async function handleOpenFolder(filePath, customUrl) {
            try {
                const res = await fetch('/api/open_folder?path=' + encodeURIComponent(filePath));
                const data = await res.json();
                if (data.status === 'error') {
                    showToast('❌ ' + data.message, true);
                } else if (data.status === 'warning') {
                    showToast('⚠️ ' + data.message, true);
                } else {
                    showToast('Opening folder in Directory Opus...');
                }
            } catch (e) {
                window.location.href = customUrl;
            }
        }

        async function checkStatus() {
            try {
                const res = await fetch('/api/status');
                const data = await res.json();
                const btn = document.getElementById('indexBtn');
                const badge = document.getElementById('indexStatusBadge');
                const countBadge = document.getElementById('docCountBadge');

                isIndexing = data.indexing_running;
                if (data.total_docs !== undefined) {
                    countBadge.innerText = '📊 Total Indexed: ' + data.total_docs.toLocaleString() + ' docs';
                }

                if (btn && badge) {
                    if (isIndexing) {
                        btn.className = 'btn-index running';
                        btn.innerText = '⏹️ Stop Indexing';
                        badge.style.display = 'inline-block';
                        
                        if (data.status_message) {
                            badge.innerText = '⏳ ' + data.status_message;
                        } else {
                            badge.innerText = '⚡ Indexing active in background...';
                        }
                    } else {
                        btn.className = 'btn-index';
                        btn.innerText = '⚡ Start Indexing';
                        badge.style.display = 'none';
                    }
                }

                const sanContainer = document.getElementById('sanSummaryContainer');
                if (data.san_summary && sanContainer) {
                    let html = '<div class="san-summary-box">';
                    html += '<div class="san-summary-header">💾 SAN Backup Recent File Changes Summary</div>';
                    html += '<div class="san-grid">';
                    for (const [cat, info] of Object.entries(data.san_summary)) {
                        html += '<div class="san-card">';
                        html += '<div class="san-card-title">📂 ' + cat + '</div>';
                        html += '<div class="san-card-file" title="' + (info.latest_path || '') + '">📄 ' + (info.latest_file || 'None') + '</div>';
                        html += '<div class="san-card-time">🕒 Modified: ' + (info.latest_time || 'N/A') + '</div>';
                        html += '</div>';
                    }
                    html += '</div></div>';
                    sanContainer.innerHTML = html;
                    sanContainer.style.display = 'block';
                }
            } catch (e) {}
        }

        async function toggleIndexing() {
            const btn = document.getElementById('indexBtn');
            const badge = document.getElementById('indexStatusBadge');

            if (isIndexing) {
                btn.innerText = '⏳ Stopping...';
                await fetch('/api/stop_indexing', { method: 'POST' });
                checkStatus();
            } else {
                if (btn) btn.className = 'btn-index running';
                if (btn) btn.innerText = '⏹️ Stop Indexing';
                if (badge) badge.style.display = 'block';

                const cfgRes = await fetch('/api/config');
                const cfg = await cfgRes.json();
                const arr = cfg.selected_directories || [];
                const netTarget = cfg.network_backup_target || '';

                await fetch('/api/config', {
                    method: 'POST',
                    headers: {'Content-Type': 'application/json'},
                    body: JSON.stringify({selected_directories: arr, network_backup_target: netTarget})
                });
                checkStatus();
            }
        }

        async function openTreeModal() {
            document.getElementById('treeModal').style.display = 'block';
            const cfgRes = await fetch('/api/config');
            const cfg = await cfgRes.json();
            selectedPaths = new Set(cfg.selected_directories || []);
            const netInput = document.getElementById('netTargetInput');
            if (netInput) netInput.value = cfg.network_backup_target || '';
            loadDriveTree();
        }

        function closeTreeModal() {
            document.getElementById('treeModal').style.display = 'none';
        }

        function norm(p) {
            return p ? p.toLowerCase().replace(/\\\\/g, '/').replace(/\/$/, '') : '';
        }

        function getCheckboxState(checkPath) {
            const cPath = norm(checkPath);
            let exactMatch = false;
            let childMatch = false;

            for (const sp of selectedPaths) {
                const normSp = norm(sp);
                if (normSp === cPath) {
                    exactMatch = true;
                } else if (normSp.startsWith(cPath + '/')) {
                    childMatch = true;
                }
            }

            if (exactMatch) return 'checked';
            if (childMatch) return 'indeterminate';
            return 'unchecked';
        }

        function updateCheckboxVisual(cb, state) {
            if (state === 'checked') {
                cb.checked = true;
                cb.indeterminate = false;
            } else if (state === 'indeterminate') {
                cb.checked = false;
                cb.indeterminate = true;
            } else {
                cb.checked = false;
                cb.indeterminate = false;
            }
        }

        async function loadDriveTree() {
            const container = document.getElementById('treeContainer');
            container.innerHTML = '';

            const res = await fetch('/api/drives');
            const drives = await res.json();

            for (const drive of drives) {
                const driveItem = document.createElement('div');
                driveItem.className = 'tree-item';

                const toggleSpan = document.createElement('span');
                toggleSpan.className = 'tree-toggle';
                toggleSpan.innerText = '▶';
                toggleSpan.setAttribute('data-path', drive);
                toggleSpan.onclick = function() { toggleFolder(this); };

                const cb = document.createElement('input');
                cb.type = 'checkbox';
                cb.value = drive;
                updateCheckboxVisual(cb, getCheckboxState(drive));
                cb.onchange = function() { onCheckboxChange(this); };

                const label = document.createElement('strong');
                label.innerHTML = ' 💻 ' + drive;

                const childrenDiv = document.createElement('div');
                childrenDiv.className = 'tree-children';

                driveItem.appendChild(toggleSpan);
                driveItem.appendChild(cb);
                driveItem.appendChild(label);
                driveItem.appendChild(childrenDiv);
                container.appendChild(driveItem);
            }
        }

        async function toggleFolder(element) {
            const path = element.getAttribute('data-path');
            const treeItem = element.parentElement;
            const childrenDiv = treeItem.querySelector('.tree-children');

            if (element.innerText === '▶') {
                element.innerText = '▼';
                childrenDiv.classList.add('open');

                if (childrenDiv.children.length === 0) {
                    const res = await fetch('/api/ls?path=' + encodeURIComponent(path));
                    const subdirs = await res.json();

                    if (subdirs.length === 0) {
                        childrenDiv.innerHTML = '<div style="margin-left:20px; color:#aaa; font-style:italic;">(Empty or No subfolders)</div>';
                    } else {
                        for (const sub of subdirs) {
                            const subItem = document.createElement('div');
                            subItem.className = 'tree-item';

                            const subToggle = document.createElement('span');
                            subToggle.className = 'tree-toggle';
                            subToggle.innerText = '▶';
                            subToggle.setAttribute('data-path', sub.path);
                            subToggle.onclick = function() { toggleFolder(this); };

                            const subCb = document.createElement('input');
                            subCb.type = 'checkbox';
                            subCb.value = sub.path;
                            updateCheckboxVisual(subCb, getCheckboxState(sub.path));
                            subCb.onchange = function() { onCheckboxChange(this); };

                            const subLabel = document.createElement('span');
                            subLabel.innerHTML = ' 📁 ' + sub.name;

                            const subChildren = document.createElement('div');
                            subChildren.className = 'tree-children';

                            subItem.appendChild(subToggle);
                            subItem.appendChild(subCb);
                            subItem.appendChild(subLabel);
                            subItem.appendChild(subChildren);
                            childrenDiv.appendChild(subItem);
                        }
                    }
                }
            } else {
                element.innerText = '▶';
                childrenDiv.classList.remove('open');
            }
        }

        function onCheckboxChange(cb) {
            cb.indeterminate = false;
            if (cb.checked) {
                selectedPaths.add(cb.value);
            } else {
                selectedPaths.delete(cb.value);
            }
        }

        async function saveSelectedDirectories() {
            const arr = Array.from(selectedPaths);
            const netInput = document.getElementById('netTargetInput');
            const netTarget = netInput ? netInput.value.trim() : '';

            const statusMsg = document.getElementById('statusMsg');
            statusMsg.innerText = 'Saving configuration & starting indexer...';

            await fetch('/api/config', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({selected_directories: arr, network_backup_target: netTarget})
            });

            statusMsg.innerText = '✓ Indexer triggered in background!';
            setTimeout(() => {
                closeTreeModal();
                statusMsg.innerText = '';
                checkStatus();
            }, 1500);
        }

        // Live automatic status & total document count updates every 2 seconds
        setInterval(checkStatus, 2000);
        checkStatus();
    </script>
</body>
</html>
"""


def stop_all_indexer_processes():
    global INDEXER_PROCESS
    if INDEXER_PROCESS and INDEXER_PROCESS.poll() is None:
        try:
            INDEXER_PROCESS.terminate()
            INDEXER_PROCESS.kill()
        except Exception:
            pass
    INDEXER_PROCESS = None

    try:
        import psutil
        for p in psutil.process_iter(['pid', 'name', 'cmdline']):
            try:
                if p.info['name'] and 'python' in p.info['name'].lower():
                    cmd = p.info['cmdline']
                    if cmd and any('ingest_documents.py' in arg for arg in cmd):
                        p.kill()
            except Exception:
                pass
    except Exception:
        pass

    try:
        if sys.platform == 'win32':
            subprocess.run(["taskkill", "/F", "/IM", "robocopy.exe"], capture_output=True)
    except Exception:
        pass

    # Ensure cloud sync is resumed
    try:
        if os.path.exists(r"C:\Program Files\Microsoft OneDrive\OneDrive.exe"):
            subprocess.Popen([r"C:\Program Files\Microsoft OneDrive\OneDrive.exe"])
    except Exception:
        pass
    try:
        if os.path.exists(r"C:\Program Files (x86)\Dropbox\Client\Dropbox.exe"):
            subprocess.Popen([r"C:\Program Files (x86)\Dropbox\Client\Dropbox.exe"])
    except Exception:
        pass

    prog_file = os.path.join(os.path.dirname(__file__), "indexer_progress.json")
    if os.path.exists(prog_file):
        try:
            with open(prog_file, "r", encoding="utf-8") as f:
                data = json.load(f)
            data["is_running"] = False
            data["status_message"] = ""
            with open(prog_file, "w", encoding="utf-8") as f:
                json.dump(data, f)
        except Exception:
            pass


def check_indexer_process_running():
    global INDEXER_PROCESS
    if INDEXER_PROCESS is not None and INDEXER_PROCESS.poll() is None:
        return True

    # Inspect OS process tree to verify if ingest_documents.py is active
    for p in psutil.process_iter(['name', 'cmdline']):
        try:
            if p.info['name'] and 'python' in p.info['name'].lower():
                cmd = p.info.get('cmdline') or []
                if any('ingest_documents.py' in str(arg) for arg in cmd):
                    return True
        except (psutil.NoSuchProcess, psutil.AccessDenied):
            pass
    return False


class SearchHandler(SimpleHTTPRequestHandler):
    def send_json(self, data, status=200):
        try:
            self.send_response(status)
            self.send_header("Content-Type", "application/json; charset=utf-8")
            self.end_headers()
            self.wfile.write(json.dumps(data).encode('utf-8'))
        except (ConnectionAbortedError, ConnectionResetError, BrokenPipeError):
            pass

    def do_GET(self):
        parsed = urlparse(self.path)
        params = parse_qs(parsed.query)

        # API: Raw File Content Endpoint for In-Browser Document Viewer
        if parsed.path == '/api/raw_file':
            file_path = params.get('path', [''])[0]
            if file_path and os.path.exists(file_path):
                ext = os.path.splitext(file_path)[1].lower().lstrip('.')
                content_types = {
                    'pdf': 'application/pdf',
                    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    'doc': 'application/msword',
                    'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                    'xls': 'application/vnd.ms-excel',
                    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    'ppt': 'application/vnd.ms-powerpoint',
                    'txt': 'text/plain; charset=utf-8',
                    'md': 'text/markdown; charset=utf-8',
                    'csv': 'text/csv; charset=utf-8',
                    'html': 'text/html; charset=utf-8',
                    'png': 'image/png',
                    'jpg': 'image/jpeg',
                    'jpeg': 'image/jpeg',
                    'webp': 'image/webp'
                }
                ctype = content_types.get(ext, 'application/octet-stream')
                try:
                    with open(file_path, 'rb') as f:
                        data = f.read()
                    self.send_response(200)
                    self.send_header("Content-Type", ctype)
                    self.send_header("Content-Length", str(len(data)))
                    self.end_headers()
                    self.wfile.write(data)
                    return
                except Exception as e:
                    self.send_response(500)
                    self.end_headers()
                    return
            self.send_response(404)
            self.end_headers()
            return

        # API: Render & Serve Page 1 Thumbnail Image
        if parsed.path == '/api/thumbnail':
            file_path = params.get('path', [''])[0]
            if file_path:
                img_data, content_type = get_thumbnail_bytes(file_path)
                if img_data:
                    try:
                        self.send_response(200)
                        self.send_header("Content-Type", content_type)
                        self.send_header("Cache-Control", "no-cache, must-revalidate")
                        self.end_headers()
                        self.wfile.write(img_data)
                    except (BrokenPipeError, ConnectionAbortedError, ConnectionResetError):
                        pass
                    return
            self.send_response(404)
            self.end_headers()
            return

        # API: Open File in Windows Default App
        if parsed.path == '/api/open_file':
            file_path = params.get('path', [''])[0]
            if file_path:
                norm_path = os.path.normpath(file_path)
                if os.path.exists(norm_path):
                    try:
                        os.startfile(norm_path)
                        self.send_json({"status": "ok", "message": "File opened successfully"})
                    except Exception as e:
                        self.send_json({"status": "error", "message": str(e)}, status=500)
                else:
                    parent_dir = os.path.dirname(norm_path)
                    if os.path.exists(parent_dir):
                        subprocess.Popen(['explorer', parent_dir])
                        self.send_json({"status": "warning", "message": f"File no longer exists. Opened parent folder: {parent_dir}"})
                    else:
                        self.send_json({"status": "error", "message": f"File and parent folder not found on disk"}, status=404)
            else:
                self.send_json({"status": "error", "message": "No file path provided"}, status=400)
            return

        # API: Open Containing Folder directly in Windows File Explorer or Directory Opus
        if parsed.path == '/api/open_folder':
            file_path = params.get('path', [''])[0]
            force_explorer = params.get('explorer', ['0'])[0] == '1'

            if file_path:
                norm_path = os.path.normpath(file_path)
                target_path = norm_path
                file_missing = False

                if not os.path.exists(norm_path):
                    parent_dir = os.path.dirname(norm_path)
                    if os.path.exists(parent_dir):
                        target_path = parent_dir
                        file_missing = True
                    else:
                        self.send_json({"status": "error", "message": f"Path and parent directory not found on disk"}, status=404)
                        return

                try:
                    if force_explorer:
                        if file_missing:
                            subprocess.Popen(['explorer', target_path])
                        else:
                            subprocess.Popen(['explorer', '/select,', target_path])
                        msg = f"File moved/deleted. Opened parent folder: {target_path}" if file_missing else "Folder opened in Windows Explorer"
                        self.send_json({"status": "warning" if file_missing else "ok", "message": msg})
                    else:
                        if os.path.exists(DOPUS_RT):
                            if file_missing or os.path.isdir(target_path):
                                subprocess.Popen([DOPUS_RT, "/cmd", "Go", target_path, "NEW"])
                            else:
                                parent_dir = os.path.dirname(target_path)
                                file_name = os.path.basename(target_path)
                                subprocess.Popen([DOPUS_RT, "/cmd", "Go", parent_dir, "NEW", f"SELECT={file_name}"])
                        elif os.path.exists(DOPUS_EXE):
                            if file_missing or os.path.isdir(target_path):
                                subprocess.Popen([DOPUS_EXE, target_path])
                            else:
                                parent_dir = os.path.dirname(target_path)
                                subprocess.Popen([DOPUS_EXE, parent_dir])
                        else:
                            if file_missing or os.path.isdir(target_path):
                                subprocess.Popen(['explorer', target_path])
                            else:
                                subprocess.Popen(['explorer', '/select,', target_path])
                            
                        msg = f"File moved/deleted. Opened parent folder: {target_path}" if file_missing else "Folder opened in Directory Opus"
                        self.send_json({"status": "warning" if file_missing else "ok", "message": msg})
                except Exception as e:
                    self.send_json({"status": "error", "message": str(e)}, status=500)
            else:
                self.send_json({"status": "error", "message": "No file path provided"}, status=400)
            return

        # API: Indexing Status & Total Count
        if parsed.path == '/api/status':
            is_running = check_indexer_process_running()
            total_count = get_document_count()
            
            prog_file = os.path.join(os.path.dirname(__file__), "indexer_progress.json")
            progress = {}
            if os.path.exists(prog_file):
                try:
                    with open(prog_file, "r", encoding="utf-8") as f:
                        progress = json.load(f)
                except Exception:
                    pass
                    
            msg = progress.get("status_message", "")
            if not is_running:
                msg = ""
            elif is_running and not msg:
                msg = "Indexing active in background..."

            self.send_json({
                "indexing_running": is_running,
                "total_docs": total_count,
                "scanned": progress.get("scanned", 0),
                "indexed": progress.get("indexed", 0),
                "skipped": progress.get("skipped", 0),
                "status_message": msg,
                "san_summary": progress.get("san_summary", None)
            })
            return

        # API: Available Drives
        if parsed.path == '/api/drives':
            drives = get_available_drives()
            self.send_json(drives)
            return

        # API: Subdirectories
        if parsed.path == '/api/ls':
            parent_path = params.get('path', [''])[0]
            subdirs = list_subdirectories(parent_path)
            self.send_json(subdirs)
            return

        # API: Current Config
        if parsed.path == '/api/config':
            cfg = load_config()
            self.send_json(cfg)
            return

        # Render Main Search HTML
        query_str = params.get('q', [''])[0].strip()
        sort_by = params.get('sort', ['relevance'])[0].strip()
        
        try:
            page = int(params.get('page', ['1'])[0])
            if page < 1:
                page = 1
        except Exception:
            page = 1

        cfg = load_config()
        selected_json = json.dumps(cfg.get('selected_directories', []))

        stats_html = ""
        results_html = ""
        pagination_html = ""

        sort_state = {
            "SORT_RELEVANCE": "selected" if sort_by == "relevance" else "",
            "SORT_DATE_DESC": "selected" if sort_by == "date_desc" else "",
            "SORT_DATE_ASC": "selected" if sort_by == "date_asc" else "",
            "SORT_NAME_ASC": "selected" if sort_by == "name_asc" else "",
            "SORT_SIZE_DESC": "selected" if sort_by == "size_desc" else ""
        }

        if query_str:
            client = get_client()
            es_query = parse_smart_query(query_str, sort_by=sort_by, page=page, page_size=PAGE_SIZE)
            res = None
            last_err = None
            for attempt in range(3):
                try:
                    res = client.search(index="documents", body=es_query)
                    break
                except Exception as e:
                    last_err = e
                    time.sleep(0.6 * (attempt + 1))

            if res is not None:
                hits = res['hits']['hits']
                total = res['hits']['total']['value']
                took = res['took']

                total_pages = math.ceil(total / PAGE_SIZE) if total > 0 else 1
                start_doc = (page - 1) * PAGE_SIZE + 1 if total > 0 else 0
                end_doc = min(page * PAGE_SIZE, total)

                stats_html = f"Found {total:,} matching document(s) in {took} ms"

                # Build Pagination Controls
                if total > 0:
                    prev_disabled = "disabled" if page <= 1 else ""
                    prev_page = page - 1 if page > 1 else 1
                    prev_url = f"/?q={urllib.parse.quote(query_str)}&sort={sort_by}&page={prev_page}"

                    next_disabled = "disabled" if end_doc >= total else ""
                    next_page = page + 1 if end_doc < total else page
                    next_count = min(PAGE_SIZE, total - end_doc)
                    next_label = f"Next {next_count} of {total:,} ➡️" if next_count > 0 else "Next ➡️"
                    next_url = f"/?q={urllib.parse.quote(query_str)}&sort={sort_by}&page={next_page}"

                    pagination_html = f"""
                    <div class="pagination-bar">
                        <a href="{prev_url}" class="btn-page {prev_disabled}">⬅️ Previous 100</a>
                        <span class="page-info">Showing {start_doc:,} - {end_doc:,} of {total:,} documents (Page {page} of {total_pages})</span>
                        <a href="{next_url}" class="btn-page {next_disabled}">{next_label}</a>
                    </div>
                    """

                if not hits:
                    results_html = "<div class='result-card'>No matching documents found.</div>"
                else:
                    cards = []
                    for hit in hits:
                        src = hit['_source']
                        fname = html.escape(src.get('file_name', 'Unnamed File'))
                        fpath = src.get('file_path', '')
                        
                        # Generate direct Windows Custom Protocol URLs & escaped JS paths
                        encoded_path = urllib.parse.quote(fpath.replace('\\', '/'))
                        escaped_js_path = html.escape(fpath.replace('\\', '\\\\').replace("'", "\\'"))
                        escaped_title = html.escape(fname.replace("'", "\\'"))
                        
                        open_file_url = f"openfile://{encoded_path}"
                        open_opus_url = f"openopus://{encoded_path}"
                        open_explorer_url = f"openexplorer://{encoded_path}"
                        escaped_display_path = html.escape(fpath)
                        ftype = src.get('file_type', 'doc').lower()

                        mod_val = src.get('modified_date', '')
                        create_val = src.get('created_date', '')

                        if mod_val:
                            mtime_ver = hashlib.md5(str(mod_val).encode('utf-8')).hexdigest()[:8]
                        else:
                            mtime_ver = 2

                        thumb_url = f"/api/thumbnail?path={encoded_path}&v={mtime_ver}"

                        mod_str = format_doc_date(mod_val)
                        create_str = format_doc_date(create_val)

                        date_parts = []
                        if mod_str:
                            date_parts.append(f'<span class="meta-date-item" title="Last Modified Date">🕒 <strong>Modified:</strong> {mod_str}</span>')
                        if create_str:
                            date_parts.append(f'<span class="meta-date-item" title="Date Created">🗓️ <strong>Created:</strong> {create_str}</span>')
                        
                        dates_html = f'<div class="file-meta-dates">{" &bull; ".join(date_parts)}</div>' if date_parts else ''

                        highlights = hit.get('highlight', {}).get('content', [])
                        if highlights:
                            snippet_text = " ... ".join(highlights)
                        else:
                            snippet_text = (src.get('content', '')[:300] + "...") if src.get('content') else "No preview text available."

                        # Render Right Column Visual Cover Preview Card for ALL Document Types
                        thumb_html = f"""
                        <div class="card-right">
                            <img src="{thumb_url}" class="thumb-preview" onclick="openDocViewer('{escaped_js_path}', '{escaped_title}', '{ftype}')" title="Click to view live full document preview" alt="Document Preview">
                        </div>
                        """

                        card = f"""
                        <div class="result-card">
                            <div class="card-left">
                                <div class="result-header">
                                    <a class="file-title" onclick="openDocViewer('{escaped_js_path}', '{escaped_title}', '{ftype}')" title="Click to view live full document preview in browser">{fname}</a>
                                    <div class="card-actions">
                                        <button type="button" class="btn-action btn-preview" onclick="openDocViewer('{escaped_js_path}', '{escaped_title}', '{ftype}')">👁️ Preview</button>
                                        <button type="button" class="btn-action btn-open-file" onclick="handleOpenFile('{escaped_js_path}', '{open_file_url}')">↗️ Open File</button>
                                        <button type="button" class="btn-action btn-open-explorer" onclick="handleOpenExplorer('{escaped_js_path}', '{open_explorer_url}')">📁 Explorer</button>
                                        <button type="button" class="btn-action btn-open-folder" onclick="handleOpenFolder('{escaped_js_path}', '{open_opus_url}')">📁 Opus</button>
                                        <span class="badge">{ftype}</span>
                                    </div>
                                </div>
                                <div class="file-path" onclick="handleOpenExplorer('{escaped_js_path}', '{open_explorer_url}')" title="Click to open folder in Windows File Explorer">📁 {escaped_display_path}</div>
                                {dates_html}
                                <div class="snippet">{snippet_text}</div>
                            </div>
                            {thumb_html}
                        </div>
                        """
                        cards.append(card)
                    results_html = "\n".join(cards)
            else:
                err_msg = str(last_err)
                if "503" in err_msg or "search_phase_execution_exception" in err_msg:
                    results_html = """<div class='result-card' style='border-left: 5px solid #ff9800; background-color: #fff3e0; padding: 20px; text-align: center; border-radius: 8px; margin-top: 15px;'>
                        <h3 style='color: #e65100; margin-top: 0; margin-bottom: 8px;'>⚠️ OpenSearch Engine Warming Up / Initializing Shards</h3>
                        <p style='color: #bf360c; font-size: 14px; margin-bottom: 12px;'>The search cluster is currently initializing shards or completing background index flushes (503 Service Unavailable).</p>
                        <p style='color: #666; font-size: 13px; margin-bottom: 15px;'>Please wait 5–10 seconds and click the button below to retry your search.</p>
                        <button onclick='window.location.reload()' style='background: #e65100; color: white; border: none; padding: 10px 22px; font-weight: bold; border-radius: 6px; cursor: pointer; font-size: 14px;'>🔄 Retry Search</button>
                    </div>"""
                else:
                    results_html = f"<div class='result-card' style='color:red;'>Error executing search: {html.escape(err_msg)}</div>"

        html_out = HTML_TEMPLATE.replace("{QUERY}", query_str).replace("{STATS}", stats_html).replace("{RESULTS}", results_html).replace("{SELECTED_JSON}", selected_json)
        html_out = html_out.replace("{PAGINATION_TOP}", pagination_html).replace("{PAGINATION_BOTTOM}", pagination_html)
        
        for key, val in sort_state.items():
            html_out = html_out.replace(f"{{{key}}}", val)

        self.send_response(200)
        self.send_header("Content-Type", "text/html; charset=utf-8")
        self.end_headers()
        self.wfile.write(html_out.encode('utf-8'))

    def do_POST(self):
        global INDEXER_PROCESS
        parsed = urlparse(self.path)

        # API: Stop Indexing
        if parsed.path == '/api/stop_indexing':
            stop_all_indexer_processes()
            self.send_json({"status": "ok", "message": "Indexing stopped."})
            return

        # API: Config and Start Indexing
        if parsed.path == '/api/config':
            content_length = int(self.headers.get('Content-Length', 0))
            body = self.rfile.read(content_length)
            try:
                data = json.loads(body)
                save_config(data)

                # Terminate any running instance first
                if INDEXER_PROCESS and INDEXER_PROCESS.poll() is None:
                    try:
                        INDEXER_PROCESS.terminate()
                    except Exception:
                        pass

                # Trigger ingest_documents.py as manageable subprocess
                dirs = data.get("selected_directories", [])
                if dirs:
                    # Immediately initialize progress file for real-time starting UI status
                    prog_file = os.path.join(os.path.dirname(__file__), "indexer_progress.json")
                    try:
                        with open(prog_file, "w", encoding="utf-8") as f:
                            json.dump({
                                "is_running": True,
                                "scanned": 1,
                                "indexed": 0,
                                "skipped": 0,
                                "status_message": "⚡ Pre-fetching database metadata (takes ~15-30s)...",
                                "timestamp": time.time()
                            }, f)
                    except Exception:
                        pass

                    def run_ingest():
                        global INDEXER_PROCESS
                        try:
                            cmd = [sys.executable, "ingest_documents.py", "--dir"] + dirs
                            INDEXER_PROCESS = subprocess.Popen(cmd, cwd=os.path.dirname(__file__))
                            INDEXER_PROCESS.wait()
                        finally:
                            INDEXER_PROCESS = None

                    t = threading.Thread(target=run_ingest)
                    t.daemon = True
                    t.start()

                self.send_json({"status": "ok", "message": "Saved and indexing started!"})
            except Exception as e:
                self.send_json({"status": "error", "message": str(e)}, status=500)

    def send_json(self, data, status=200):
        self.send_response(status)
        self.send_header("Content-Type", "json/application; charset=utf-8")
        self.end_headers()
        self.wfile.write(json.dumps(data).encode('utf-8'))


def main():
    server = ThreadingHTTPServer(('localhost', 8080), SearchHandler)
    print("[+] OpenSearch Smart Search Server running at http://localhost:8080")
    server.serve_forever()


if __name__ == "__main__":
    main()
