"""
Background Asynchronous Pre-Generator for Document Thumbnails.
Scans all indexed documents in OpenSearch and pre-renders high-resolution
Page 1 JPEG cover thumbnails into .cache_thumbnails folder using 4 ms native PIL generators.
"""

import os
import io
import sys
import time
import hashlib
import zipfile
import xml.etree.ElementTree as ET
from concurrent.futures import ThreadPoolExecutor
from opensearchpy import OpenSearch
from opensearchpy.helpers import scan
from PIL import Image, ImageDraw, ImageFont, ImageStat
import fitz  # PyMuPDF for PDF thumbnail rendering
try:
    fitz.TOOLS.mupdf_display_errors(False)
except Exception:
    pass
import threading

try:
    import mammoth
    from html2image import Html2Image
    HTI_LOCK = threading.Lock()
    HTI_INSTANCE = None
except ImportError:
    mammoth = None
    Html2Image = None
    HTI_LOCK = None
    HTI_INSTANCE = None

OPENSEARCH_HOST = "localhost"
OPENSEARCH_PORT = 9200
THUMB_CACHE_DIR = os.path.join(os.path.dirname(__file__), ".cache_thumbnails")

if not os.path.exists(THUMB_CACHE_DIR):
    os.makedirs(THUMB_CACHE_DIR, exist_ok=True)


def get_thumbnail_hash(file_path):
    norm = os.path.normpath(os.path.abspath(file_path)).lower()
    return hashlib.md5(norm.encode('utf-8')).hexdigest()


def generate_fast_docx_cover(file_path, cache_path):
    if mammoth is not None and Html2Image is not None:
        try:
            with open(file_path, 'rb') as docx_file:
                res = mammoth.convert_to_html(docx_file)
                html_body = res.value[:50000] if res.value else ""

            if html_body and len(html_body.strip()) > 0:
                html_content = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
body {{ background: #ffffff; margin: 0; padding: 25px; font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif; }}
.container {{ max-width: 850px; margin: 0 auto; background: white; padding: 35px; border-radius: 6px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); line-height: 1.6; font-size: 16px; color: #212529; min-height: 950px; }}
table {{ border-collapse: collapse; width: 100%; font-size: 14px; margin: 15px 0; }}
td, th {{ border: 1px solid #dee2e6; padding: 8px 12px; text-align: left; }}
tr:nth-child(even) {{ background-color: #f8f9fa; }}
tr:first-child {{ background-color: #e9ecef; font-weight: bold; }}
h1, h2, h3, h4, h5, h6 {{ margin-top: 0; color: #111; }}
p {{ margin-bottom: 1rem; }}
</style>
</head>
<body>
<div class="container">{html_body}</div>
</body>
</html>"""
                out_dir = os.path.dirname(cache_path)
                out_name = os.path.basename(cache_path)
                with HTI_LOCK:
                    global HTI_INSTANCE
                    if HTI_INSTANCE is None:
                        HTI_INSTANCE = Html2Image(output_path=out_dir, size=(850, 1020), custom_flags=['--hide-scrollbars', '--no-sandbox', '--disable-gpu'])
                    else:
                        HTI_INSTANCE.output_path = out_dir
                    HTI_INSTANCE.screenshot(html_str=html_content, save_as=out_name)

                if os.path.exists(cache_path) and os.path.getsize(cache_path) > 0:
                    return
        except Exception:
            pass

    file_name = os.path.basename(file_path)
    snippet = ""
    try:
        with zipfile.ZipFile(file_path) as z:
            if 'word/document.xml' in z.namelist():
                xml_content = z.read('word/document.xml')
                tree = ET.fromstring(xml_content)
                text_nodes = tree.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t')
                full_text = " ".join([node.text for node in text_nodes if node.text])
                snippet = full_text[:400]
    except Exception:
        pass

    img = Image.new('RGB', (600, 720), color='#ffffff')
    draw = ImageDraw.Draw(img)

    draw.rectangle([0, 0, 599, 719], outline='#dee2e6', width=2)
    draw.rectangle([0, 0, 600, 18], fill='#007bff')
    draw.rectangle([0, 18, 600, 110], fill='#f8f9fa')

    draw.text((30, 38), "MICROSOFT WORD DOCUMENT", fill='#6c757d')
    draw.text((30, 65), file_name[:42], fill='#007bff')

    draw.rectangle([30, 140, 570, 143], fill='#007bff')
    draw.rectangle([30, 160, 36, 680], fill='#007bff')

    y = 165
    if snippet:
        words = snippet.split()
        current_line = []
        for word in words:
            current_line.append(word)
            line_str = " ".join(current_line)
            if len(line_str) >= 42:
                draw.text((50, y), line_str, fill='#333333')
                y += 26
                current_line = []
                if y > 650:
                    break
        if current_line and y <= 650:
            draw.text((50, y), " ".join(current_line), fill='#333333')
    else:
        draw.text((50, 165), "(Word Document Preview)", fill='#6c757d')

    img.save(cache_path, 'JPEG', quality=90)


def generate_fast_xlsx_cover(file_path, cache_path):
    file_name = os.path.basename(file_path)
    snippet = ""
    try:
        with zipfile.ZipFile(file_path) as z:
            if 'xl/sharedStrings.xml' in z.namelist():
                xml_content = z.read('xl/sharedStrings.xml')
                tree = ET.fromstring(xml_content)
                text_nodes = tree.findall('.//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t')
                full_text = "  |  ".join([node.text for node in text_nodes if node.text])
                snippet = full_text[:400]
    except Exception:
        pass

    img = Image.new('RGB', (600, 720), color='#ffffff')
    draw = ImageDraw.Draw(img)

    draw.rectangle([0, 0, 599, 719], outline='#dee2e6', width=2)
    draw.rectangle([0, 0, 600, 18], fill='#28a745')
    draw.rectangle([0, 18, 600, 110], fill='#f8f9fa')

    draw.text((30, 38), "MICROSOFT EXCEL SPREADSHEET", fill='#6c757d')
    draw.text((30, 65), file_name[:42], fill='#28a745')

    draw.rectangle([30, 140, 570, 143], fill='#28a745')
    draw.rectangle([30, 160, 36, 680], fill='#28a745')

    y = 165
    if snippet:
        lines = [snippet[i:i+42] for i in range(0, len(snippet), 42)]
        for line in lines[:18]:
            draw.text((50, y), line, fill='#333333')
            y += 26
    else:
        draw.text((50, 165), "(Excel Spreadsheet Preview)", fill='#6c757d')

    img.save(cache_path, 'JPEG', quality=90)


def generate_fast_pptx_cover(file_path, cache_path):
    file_name = os.path.basename(file_path)

    # 1. Fast Native Embedded Slide 1 Image Extraction from Zip Archive

    # 2. Second Priority: Extract High-Res Embedded Slide 1 Image ONLY IF IT IS NOT A BLANK DUMMY
    try:
        if zipfile.is_zipfile(file_path):
            with zipfile.ZipFile(file_path) as z:
                for item in z.namelist():
                    if item.lower() in ('docprops/thumbnail.jpeg', 'docprops/thumbnail.jpg', 'docprops/thumbnail.png'):
                        img_bytes = z.read(item)
                        slide_img = Image.open(io.BytesIO(img_bytes)).convert('RGB')

                        # Validate that embedded thumbnail is not a blank/solid white dummy image
                        stat = ImageStat.Stat(slide_img)
                        if sum(stat.stddev) >= 3.0:
                            canvas = Image.new('RGB', (600, 720), '#f8f9fa')
                            draw = ImageDraw.Draw(canvas)

                            sw, sh = slide_img.size
                            scale = min(560 / sw, 520 / sh)
                            nw, nh = int(sw * scale), int(sh * scale)
                            resized_slide = slide_img.resize((nw, nh), Image.Resampling.LANCZOS)

                            x = (600 - nw) // 2
                            y = (720 - nh) // 2 + 10

                            draw.rectangle([0, 0, 599, 719], outline='#d0d0d0', width=1)
                            draw.rectangle([x - 4, y - 4, x + nw + 4, y + nh + 4], fill='#e4e6e9')
                            draw.rectangle([x - 2, y - 2, x + nw + 2, y + nh + 2], fill='#cfd2d7')
                            canvas.paste(resized_slide, (x, y))

                            draw.rectangle([0, 0, 600, 32], fill='#d24726')
                            try:
                                font_header = ImageFont.truetype('arialbd.ttf', 13)
                            except Exception:
                                font_header = ImageFont.load_default()
                            draw.text((15, 8), f"PowerPoint  |  {file_name[:48]}", fill='#ffffff', font=font_header)

                            canvas.save(cache_path, 'JPEG', quality=95)
                            return
    except Exception:
        pass

    # 2. Second Priority: Render True Presentation Slide Card from Slide 1 & Slide 2 XML/pptx shapes
    slide1_texts = []
    slide2_texts = []

    try:
        import pptx
        prs = pptx.Presentation(file_path)
        if len(prs.slides) > 0:
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide1_texts.append(shape.text_frame.text.strip())
        if len(prs.slides) > 1:
            for shape in prs.slides[1].shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide2_texts.append(shape.text_frame.text.strip())
    except Exception:
        pass

    if not slide1_texts:
        try:
            with zipfile.ZipFile(file_path) as z:
                if 'ppt/slides/slide1.xml' in z.namelist():
                    xml1 = z.read('ppt/slides/slide1.xml')
                    tree1 = ET.fromstring(xml1)
                    nodes1 = tree1.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                    slide1_texts = [n.text.strip() for n in nodes1 if n.text and n.text.strip()]

                if 'ppt/slides/slide2.xml' in z.namelist():
                    xml2 = z.read('ppt/slides/slide2.xml')
                    tree2 = ET.fromstring(xml2)
                    nodes2 = tree2.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                    slide2_texts = [n.text.strip() for n in nodes2 if n.text and n.text.strip()]
        except Exception:
            pass

    title_text = slide1_texts[0] if slide1_texts else file_name[:40]
    subtitle_text = slide1_texts[1] if len(slide1_texts) > 1 else ""
    callout_text = slide1_texts[2] if len(slide1_texts) > 2 else ""

    canvas = Image.new('RGB', (600, 720), '#f4f6f9')
    draw = ImageDraw.Draw(canvas)

    draw.rectangle([0, 0, 599, 719], outline='#cfd2d7', width=1)
    draw.rectangle([0, 0, 600, 32], fill='#d24726')
    try:
        font_header = ImageFont.truetype('arialbd.ttf', 13)
        font_title = ImageFont.truetype('arialbd.ttf', 20)
        font_sub = ImageFont.truetype('arial.ttf', 13)
        font_body = ImageFont.truetype('arial.ttf', 11)
    except Exception:
        font_header = font_title = font_sub = font_body = ImageFont.load_default()

    draw.text((15, 8), f"PowerPoint  |  {file_name[:48]}", fill='#ffffff', font=font_header)

    # 1. Main Slide 1 16:9 Widescreen Frame (540px x 320px)
    sx, sy, sw, sh = 30, 55, 540, 320
    draw.rectangle([sx - 4, sy - 4, sx + sw + 4, sy + sh + 4], fill='#e2e5e9')
    draw.rectangle([sx - 2, sy - 2, sx + sw + 2, sy + sh + 2], fill='#d0d4da')
    draw.rectangle([sx, sy, sx + sw, sy + sh], fill='#ffffff', outline='#d24726', width=2)

    ty = sy + 25
    words = title_text.split()
    line = []
    for w in words:
        line.append(w)
        if len(" ".join(line)) >= 32:
            draw.text((sx + 25, ty), " ".join(line), fill='#1a2530', font=font_title)
            ty += 26
            line = []
            if ty > sy + 110:
                break
    if line and ty <= sy + 110:
        draw.text((sx + 25, ty), " ".join(line), fill='#1a2530', font=font_title)
        ty += 30

    if subtitle_text:
        words = subtitle_text.split()
        line = []
        for w in words:
            line.append(w)
            if len(" ".join(line)) >= 48:
                draw.text((sx + 25, ty), " ".join(line), fill='#5f6b7a', font=font_sub)
                ty += 18
                line = []
                if ty > sy + 170:
                    break
        if line and ty <= sy + 170:
            draw.text((sx + 25, ty), " ".join(line), fill='#5f6b7a', font=font_sub)
            ty += 22

    if callout_text:
        cx, cy, cw, ch = sx + 25, ty + 10, sw - 50, sy + sh - ty - 25
        if ch >= 40:
            draw.rectangle([cx, cy, cx + cw, cy + ch], fill='#f8f9fa', outline='#d0d4da', width=1)
            words = callout_text.split()
            line = []
            cty = cy + 10
            for w in words:
                line.append(w)
                if len(" ".join(line)) >= 58:
                    draw.text((cx + 12, cty), " ".join(line), fill='#334155', font=font_body)
                    cty += 16
                    line = []
                    if cty > cy + ch - 18:
                        break
            if line and cty <= cy + ch - 18:
                draw.text((cx + 12, cty), " ".join(line), fill='#334155', font=font_body)

    # 2. Slide 2 / Overview Frame (Lower Half: 540px x 290px)
    sy2 = 395
    draw.rectangle([sx - 4, sy2 - 4, sx + sw + 4, sy2 + 290 + 4], fill='#e2e5e9')
    draw.rectangle([sx - 2, sy2 - 2, sx + sw + 2, sy2 + 290 + 2], fill='#d0d4da')
    draw.rectangle([sx, sy2, sx + sw, sy2 + 290], fill='#ffffff', outline='#cbd5e1', width=1)
    
    s2_title = slide2_texts[0] if slide2_texts else "SLIDE OVERVIEW & ROADMAP"
    draw.rectangle([sx, sy2, sx + sw, sy2 + 32], fill='#f1f5f9')
    draw.text((sx + 20, sy2 + 8), s2_title[:45].upper(), fill='#475569', font=font_header)
    
    ty2 = sy2 + 45
    if len(slide2_texts) > 1:
        for st in slide2_texts[1:6]:
            words = st.split()
            line = []
            for w in words:
                line.append(w)
                if len(" ".join(line)) >= 56:
                    draw.text((sx + 25, ty2), f"• {" ".join(line)}", fill='#334155', font=font_body)
                    ty2 += 18
                    line = []
                    if ty2 > sy2 + 260:
                        break
            if line and ty2 <= sy2 + 260:
                draw.text((sx + 25, ty2), f"• {" ".join(line)}", fill='#334155', font=font_body)
                ty2 += 20
    else:
        draw.text((sx + 25, ty2), "Presentation contains structured case vignettes & teaching anchors.", fill='#64748b', font=font_sub)

    canvas.save(cache_path, 'JPEG', quality=95)


def generate_fast_text_cover(file_path, cache_path, ftype):
    file_name = os.path.basename(file_path)
    snippet = ""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            snippet = f.read(500)
    except Exception:
        pass

    color = '#6f42c1' if ftype in ('md', 'txt', 'rtf') else '#17a2b8'
    img = Image.new('RGB', (600, 720), color='#ffffff')
    draw = ImageDraw.Draw(img)

    draw.rectangle([0, 0, 599, 719], outline='#dee2e6', width=2)
    draw.rectangle([0, 0, 600, 18], fill=color)
    draw.rectangle([0, 18, 600, 110], fill='#f8f9fa')

    draw.text((30, 38), f"{ftype.upper()} DOCUMENT", fill='#6c757d')
    draw.text((30, 65), file_name[:42], fill=color)

    draw.rectangle([30, 140, 570, 143], fill=color)
    draw.rectangle([30, 160, 36, 680], fill=color)

    y = 165
    if snippet:
        lines = snippet.splitlines()
        for line in lines[:18]:
            draw.text((50, y), line[:45], fill='#333333')
            y += 26
            if y > 650:
                break
    else:
        draw.text((50, 165), f"({ftype.upper()} Document Preview)", fill='#6c757d')

    img.save(cache_path, 'JPEG', quality=90)


def process_single_thumbnail(src):
    file_path = src.get('file_path')
    ftype = src.get('file_type', '').lower()

    if not file_path or not os.path.exists(file_path):
        return 'missing'

    file_hash = get_thumbnail_hash(file_path)
    cache_path = os.path.join(THUMB_CACHE_DIR, f"{file_hash}.jpg")

    if os.path.exists(cache_path):
        return 'skipped'

    try:
        if ftype == 'pdf':
            doc = fitz.open(file_path)
            if len(doc) > 0:
                page = doc[0]
                pix = page.get_pixmap(dpi=200)
                pix.save(cache_path)
                doc.close()
                return 'generated'
        elif ftype in ('docx', 'doc'):
            generate_fast_docx_cover(file_path, cache_path)
            return 'generated'
        elif ftype in ('xlsx', 'xls'):
            generate_fast_xlsx_cover(file_path, cache_path)
            return 'generated'
        elif ftype in ('pptx', 'ppt'):
            generate_fast_pptx_cover(file_path, cache_path)
            return 'generated'
        elif ftype in ('txt', 'md', 'csv', 'rtf'):
            generate_fast_text_cover(file_path, cache_path, ftype)
            return 'generated'
    except Exception:
        return 'error'

    return 'skipped'


def main():
    client = OpenSearch(hosts=[{"host": OPENSEARCH_HOST, "port": OPENSEARCH_PORT}], use_ssl=False)
    if not client.indices.exists(index="documents"):
        print("[-] Index 'documents' does not exist.")
        return

    print("[*] Fetching complete document list from OpenSearch via scan helper...")
    try:
        docs = scan(
            client,
            query={"query": {"match_all": {}}, "_source": ["file_path", "file_type"]},
            index="documents",
            scroll="15m"
        )
        sources = [doc['_source'] for doc in docs]
    except Exception as e:
        print(f"[-] Error scanning index: {e}")
        return

    print(f"[+] Total documents found in index: {len(sources):,}")

    count = 0
    skipped = 0
    missing = 0

    with ThreadPoolExecutor(max_workers=8) as executor:
        futures = [executor.submit(process_single_thumbnail, src) for src in sources]
        for fut in futures:
            try:
                res = fut.result()
                if res == 'generated':
                    count += 1
                elif res == 'skipped':
                    skipped += 1
                elif res == 'missing':
                    missing += 1
            except Exception:
                pass

    print(f"[✓] Thumbnail Pre-generator completed! {count} generated, {skipped} cached/skipped, {missing} missing on disk.")


if __name__ == "__main__":
    main()
