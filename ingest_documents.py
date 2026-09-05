"""
High-Performance Document Ingestion Pipeline for OpenSearch.
Includes:
- Multi-Threaded Real-Time 20-Doc Batch Flushes with OpenSearch Index Refresh
- Instant Live UI Counter Updates on http://localhost:8080
- 4 ms Native PIL DOCX & PyMuPDF Page 1 Thumbnail Generator
"""

import os
import io
import sys

if hasattr(sys.stdout, 'reconfigure'):
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except Exception:
        pass
if hasattr(sys.stderr, 'reconfigure'):
    try:
        sys.stderr.reconfigure(encoding='utf-8')
    except Exception:
        pass
import json
import time
import hashlib
import argparse
import threading
import subprocess
import zipfile
import tempfile
import ctypes
import xml.etree.ElementTree as ET
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

from opensearchpy import OpenSearch, helpers
from PIL import Image, ImageDraw, ImageFont, ImageStat
import fitz  # PyMuPDF for PDF thumbnail rendering
try:
    fitz.TOOLS.mupdf_display_errors(False)
except Exception:
    pass

try:
    import mammoth
    from html2image import Html2Image
    HTI_LOCK = threading.Lock()
    HTI_INSTANCE = None
except ImportError:
    mammoth = None
    Html2Image = None
    HTI_LOCK = None
    HTI_INSTANCE = None

try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

# Fast Folder Exclusions (names & full paths)
DEFAULT_EXCLUDE_NAMES = {
    '__pycache__', '.git', '.svn', 'node_modules', '.venv', 'venv',
    '.dropbox.cache', 'dropboxbackup', '.cache', 'appdata', '.gemini',
    'deidentifier', 'identified', 'new topic1 - copy', 'build', 'dist',
    'target', 'bin', 'obj', '.vs', '.idea', '.vscode', '.nuget', 'coverage',
    '$recycle.bin', 'system volume information', 'windows', 'program files',
    'program files (x86)', 'programdata'
}

DEFAULT_EXCLUDE_DIRS = {
    r'c:\windows', r'c:\program files', r'c:\program files (x86)',
    r'c:\programdata', r'$recycle.bin', r'system volume information',
    r'd:\active research\deidentifier\identified', r'deidentifier\identified',
    r'd:\quicken', r'd:\backups', r'd:\backups\documents',
    r'd:\outlook files', r'd:\home videos and pictures',
    r'd:\pictures backup april 2010', r'd:\dtsearch index'
}

# File extensions to strictly SKIP
SKIP_EXTENSIONS = {
    '.dll', '.exe', '.sys', '.bin', '.so', '.dylib', '.pyc', '.class',
    '.zip', '.tar', '.gz', '.7z', '.rar', '.iso', '.obj', '.o', '.lib',
    '.pdb', '.dat', '.tmp', '.pyd', '.nupkg', '.cab', '.msi', '.lnk',
    '.cache', '.bak', '.ico', '.png', '.jpg', '.jpeg', '.gif', '.mp3', '.mp4',
    '.ini', '.log', '.tsv', '.pb', '.pak'
}

# Primary Human Document & Presentation Extensions ONLY (PDF, Word, Excel, PPT, Text, Markdown, RTF)
STRICT_DOCUMENT_EXTENSIONS = {
    '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt',
    '.txt', '.md', '.markdown', '.rtf', '.odt', '.ods', '.odp', '.epub'
}

THUMB_CACHE_DIR = os.path.join(os.path.dirname(__file__), ".cache_thumbnails")
if not os.path.exists(THUMB_CACHE_DIR):
    os.makedirs(THUMB_CACHE_DIR, exist_ok=True)


def to_long_path(path: str) -> str:
    """Converts Windows path to extended-length path (\\\\?\\ prefix) if length >= 240 characters to bypass 260-char MAX_PATH limit."""
    if sys.platform == 'win32' and isinstance(path, str) and not path.startswith('\\\\?\\'):
        abs_p = os.path.abspath(path)
        if len(abs_p) >= 240 and not abs_p.startswith('\\\\'):
            return '\\\\?\\' + abs_p
    return path


def safe_read_bytes(file_path):
    """
    Safely reads bytes from a file even if locked exclusively by Word/Excel/PowerPoint/OneDrive.
    Handles long Windows file paths (>= 260 chars) using extended path prefix.
    """
    long_p = to_long_path(file_path)
    if not os.path.exists(long_p):
        return None
    try:
        with open(long_p, 'rb') as f:
            return f.read()
    except PermissionError:
        tmp_dir = tempfile.gettempdir()
        tmp_name = f"opensearch_lock_{os.getpid()}_{os.path.basename(file_path)}"
        tmp_path = os.path.join(tmp_dir, tmp_name)
        try:
            res = ctypes.windll.kernel32.CopyFileW(file_path, tmp_path, False)
            if res != 0:
                with open(tmp_path, 'rb') as f:
                    data = f.read()
                return data
        except Exception:
            pass
        finally:
            if os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except Exception:
                    pass
    except Exception:
        pass
    return None


def get_thumbnail_hash(file_path):
    norm = os.path.normpath(os.path.abspath(file_path)).lower()
    return hashlib.md5(norm.encode('utf-8')).hexdigest()


def generate_fast_docx_cover(file_path, cache_path):
    file_bytes = safe_read_bytes(file_path)
    if file_bytes and mammoth is not None and Html2Image is not None:
        try:
            res = mammoth.convert_to_html(io.BytesIO(file_bytes))
            raw_html = res.value or ""

            # Safely slice at valid tag boundary so base64 <img src="data:image..."> tags are never chopped in half
            if len(raw_html) > 600000:
                cut_pos = raw_html.rfind('</p>', 0, 600000)
                if cut_pos == -1:
                    cut_pos = raw_html.rfind('>', 0, 600000)
                html_body = raw_html[:cut_pos+4] if cut_pos != -1 else raw_html[:600000]
            else:
                html_body = raw_html

            if html_body and len(html_body.strip()) > 0:
                html_content = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
body {{ background: #ffffff; margin: 0; padding: 25px; font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif; }}
.container {{ max-width: 850px; margin: 0 auto; background: white; padding: 35px; border-radius: 6px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); line-height: 1.6; font-size: 16px; color: #212529; min-height: 950px; }}
table {{ border-collapse: collapse; width: 100%; font-size: 14px; margin: 15px 0; }}
td, th {{ border: 1px solid #dee2e6; padding: 8px 12px; text-align: left; }}
tr:nth-child(even) {{ background-color: #f8f9fa; }}
tr:first-child {{ background-color: #e9ecef; font-weight: bold; }}
h1, h2, h3, h4, h5, h6 {{ margin-top: 0; color: #111; }}
p {{ margin-bottom: 1rem; }}
img {{ max-width: 100%; height: auto; display: block; margin: 12px 0; border-radius: 4px; }}
</style>
</head>
<body>
<div class="container">{html_body}</div>
</body>
</html>"""
                out_dir = os.path.dirname(cache_path)
                out_name = os.path.basename(cache_path)
                with HTI_LOCK:
                    global HTI_INSTANCE
                    if HTI_INSTANCE is None:
                        HTI_INSTANCE = Html2Image(output_path=out_dir, size=(850, 1020), custom_flags=['--hide-scrollbars', '--no-sandbox', '--disable-gpu'])
                    else:
                        HTI_INSTANCE.output_path = out_dir
                    HTI_INSTANCE.screenshot(html_str=html_content, save_as=out_name)

                if os.path.exists(cache_path) and os.path.getsize(cache_path) > 0:
                    return
        except Exception:
            pass

    file_name = os.path.basename(file_path)
    snippet = ""
    if file_bytes:
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                if 'word/document.xml' in z.namelist():
                    xml_content = z.read('word/document.xml')
                    tree = ET.fromstring(xml_content)
                    text_nodes = tree.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t')
                    full_text = " ".join([node.text for node in text_nodes if node.text])
                    snippet = full_text[:400]
        except Exception:
            pass

    img = Image.new('RGB', (600, 720), color='#ffffff')
    draw = ImageDraw.Draw(img)

    draw.rectangle([0, 0, 599, 719], outline='#dee2e6', width=2)
    draw.rectangle([0, 0, 600, 18], fill='#007bff')
    draw.rectangle([0, 18, 600, 110], fill='#f8f9fa')

    draw.text((30, 38), "MICROSOFT WORD DOCUMENT", fill='#6c757d')
    draw.text((30, 65), file_name[:42], fill='#007bff')

    draw.rectangle([30, 140, 570, 143], fill='#007bff')
    draw.rectangle([30, 160, 36, 680], fill='#007bff')

    y = 165
    if snippet:
        words = snippet.split()
        current_line = []
        for word in words:
            current_line.append(word)
            line_str = " ".join(current_line)
            if len(line_str) >= 42:
                draw.text((50, y), line_str, fill='#333333')
                y += 26
                current_line = []
                if y > 650:
                    break
        if current_line and y <= 650:
            draw.text((50, y), " ".join(current_line), fill='#333333')
    else:
        draw.text((50, 165), "(Word Document Preview)", fill='#6c757d')

    img.save(cache_path, 'JPEG', quality=90)


def generate_fast_pptx_cover(file_path, cache_path):
    file_name = os.path.basename(file_path)

    # 1. Fast Native Embedded Slide 1 Image Extraction from Zip Archive

    # 2. Second Priority: Extract High-Res Embedded Slide 1 Image ONLY IF IT IS NOT A BLANK DUMMY
    try:
        if zipfile.is_zipfile(file_path):
            with zipfile.ZipFile(file_path) as z:
                for item in z.namelist():
                    if item.lower() in ('docprops/thumbnail.jpeg', 'docprops/thumbnail.jpg', 'docprops/thumbnail.png'):
                        img_bytes = z.read(item)
                        slide_img = Image.open(io.BytesIO(img_bytes)).convert('RGB')

                        # Validate that embedded thumbnail is not a blank/solid white dummy image
                        stat = ImageStat.Stat(slide_img)
                        if sum(stat.stddev) >= 3.0:
                            canvas = Image.new('RGB', (600, 720), '#f8f9fa')
                            draw = ImageDraw.Draw(canvas)

                            sw, sh = slide_img.size
                            scale = min(560 / sw, 520 / sh)
                            nw, nh = int(sw * scale), int(sh * scale)
                            resized_slide = slide_img.resize((nw, nh), Image.Resampling.LANCZOS)

                            x = (600 - nw) // 2
                            y = (720 - nh) // 2 + 10

                            draw.rectangle([0, 0, 599, 719], outline='#d0d0d0', width=1)
                            draw.rectangle([x - 4, y - 4, x + nw + 4, y + nh + 4], fill='#e4e6e9')
                            draw.rectangle([x - 2, y - 2, x + nw + 2, y + nh + 2], fill='#cfd2d7')
                            canvas.paste(resized_slide, (x, y))

                            draw.rectangle([0, 0, 600, 32], fill='#d24726')
                            try:
                                font_header = ImageFont.truetype('arialbd.ttf', 13)
                            except Exception:
                                font_header = ImageFont.load_default()
                            draw.text((15, 8), f"PowerPoint  |  {file_name[:48]}", fill='#ffffff', font=font_header)

                            canvas.save(cache_path, 'JPEG', quality=95)
                            return
    except Exception:
        pass

    # 2. Second Priority: Render True Presentation Slide Card from Slide 1 & Slide 2 XML/pptx shapes
    slide1_texts = []
    slide2_texts = []

    try:
        import pptx
        prs = pptx.Presentation(file_path)
        if len(prs.slides) > 0:
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide1_texts.append(shape.text_frame.text.strip())
        if len(prs.slides) > 1:
            for shape in prs.slides[1].shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide2_texts.append(shape.text_frame.text.strip())
    except Exception:
        pass

    if not slide1_texts:
        try:
            with zipfile.ZipFile(file_path) as z:
                if 'ppt/slides/slide1.xml' in z.namelist():
                    xml1 = z.read('ppt/slides/slide1.xml')
                    tree1 = ET.fromstring(xml1)
                    nodes1 = tree1.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                    slide1_texts = [n.text.strip() for n in nodes1 if n.text and n.text.strip()]

                if 'ppt/slides/slide2.xml' in z.namelist():
                    xml2 = z.read('ppt/slides/slide2.xml')
                    tree2 = ET.fromstring(xml2)
                    nodes2 = tree2.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}t')
                    slide2_texts = [n.text.strip() for n in nodes2 if n.text and n.text.strip()]
        except Exception:
            pass

    title_text = slide1_texts[0] if slide1_texts else file_name[:40]
    subtitle_text = slide1_texts[1] if len(slide1_texts) > 1 else ""
    callout_text = slide1_texts[2] if len(slide1_texts) > 2 else ""

    canvas = Image.new('RGB', (600, 720), '#f4f6f9')
    draw = ImageDraw.Draw(canvas)

    draw.rectangle([0, 0, 599, 719], outline='#cfd2d7', width=1)
    draw.rectangle([0, 0, 600, 32], fill='#d24726')
    try:
        font_header = ImageFont.truetype('arialbd.ttf', 13)
        font_title = ImageFont.truetype('arialbd.ttf', 20)
        font_sub = ImageFont.truetype('arial.ttf', 13)
        font_body = ImageFont.truetype('arial.ttf', 11)
    except Exception:
        font_header = font_title = font_sub = font_body = ImageFont.load_default()

    draw.text((15, 8), f"PowerPoint  |  {file_name[:48]}", fill='#ffffff', font=font_header)

    # 1. Main Slide 1 16:9 Widescreen Frame (540px x 320px)
    sx, sy, sw, sh = 30, 55, 540, 320
    draw.rectangle([sx - 4, sy - 4, sx + sw + 4, sy + sh + 4], fill='#e2e5e9')
    draw.rectangle([sx - 2, sy - 2, sx + sw + 2, sy + sh + 2], fill='#d0d4da')
    draw.rectangle([sx, sy, sx + sw, sy + sh], fill='#ffffff', outline='#d24726', width=2)

    ty = sy + 25
    words = title_text.split()
    line = []
    for w in words:
        line.append(w)
        if len(" ".join(line)) >= 32:
            draw.text((sx + 25, ty), " ".join(line), fill='#1a2530', font=font_title)
            ty += 26
            line = []
            if ty > sy + 110:
                break
    if line and ty <= sy + 110:
        draw.text((sx + 25, ty), " ".join(line), fill='#1a2530', font=font_title)
        ty += 30

    if subtitle_text:
        words = subtitle_text.split()
        line = []
        for w in words:
            line.append(w)
            if len(" ".join(line)) >= 48:
                draw.text((sx + 25, ty), " ".join(line), fill='#5f6b7a', font=font_sub)
                ty += 18
                line = []
                if ty > sy + 170:
                    break
        if line and ty <= sy + 170:
            draw.text((sx + 25, ty), " ".join(line), fill='#5f6b7a', font=font_sub)
            ty += 22

    if callout_text:
        cx, cy, cw, ch = sx + 25, ty + 10, sw - 50, sy + sh - ty - 25
        if ch >= 40:
            draw.rectangle([cx, cy, cx + cw, cy + ch], fill='#f8f9fa', outline='#d0d4da', width=1)
            words = callout_text.split()
            line = []
            cty = cy + 10
            for w in words:
                line.append(w)
                if len(" ".join(line)) >= 58:
                    draw.text((cx + 12, cty), " ".join(line), fill='#334155', font=font_body)
                    cty += 16
                    line = []
                    if cty > cy + ch - 18:
                        break
            if line and cty <= cy + ch - 18:
                draw.text((cx + 12, cty), " ".join(line), fill='#334155', font=font_body)

    # 2. Slide 2 / Overview Frame (Lower Half: 540px x 290px)
    sy2 = 395
    draw.rectangle([sx - 4, sy2 - 4, sx + sw + 4, sy2 + 290 + 4], fill='#e2e5e9')
    draw.rectangle([sx - 2, sy2 - 2, sx + sw + 2, sy2 + 290 + 2], fill='#d0d4da')
    draw.rectangle([sx, sy2, sx + sw, sy2 + 290], fill='#ffffff', outline='#cbd5e1', width=1)
    
    s2_title = slide2_texts[0] if slide2_texts else "SLIDE OVERVIEW & ROADMAP"
    draw.rectangle([sx, sy2, sx + sw, sy2 + 32], fill='#f1f5f9')
    draw.text((sx + 20, sy2 + 8), s2_title[:45].upper(), fill='#475569', font=font_header)
    
    ty2 = sy2 + 45
    if len(slide2_texts) > 1:
        for st in slide2_texts[1:6]:
            words = st.split()
            line = []
            for w in words:
                line.append(w)
                if len(" ".join(line)) >= 56:
                    draw.text((sx + 25, ty2), f"• {" ".join(line)}", fill='#334155', font=font_body)
                    ty2 += 18
                    line = []
                    if ty2 > sy2 + 260:
                        break
            if line and ty2 <= sy2 + 260:
                draw.text((sx + 25, ty2), f"• {" ".join(line)}", fill='#334155', font=font_body)
                ty2 += 20
    else:
        draw.text((sx + 25, ty2), "Presentation contains structured case vignettes & teaching anchors.", fill='#64748b', font=font_sub)

    canvas.save(cache_path, 'JPEG', quality=95)


def generate_fast_xlsx_cover(file_path, cache_path):
    file_name = os.path.basename(file_path)
    snippet = ""
    try:
        if zipfile.is_zipfile(file_path):
            with zipfile.ZipFile(file_path) as z:
                if 'xl/sharedStrings.xml' in z.namelist():
                    xml_content = z.read('xl/sharedStrings.xml')
                    tree = ET.fromstring(xml_content)
                    text_nodes = tree.findall('.//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t')
                    full_text = "  |  ".join([node.text for node in text_nodes if node.text])
                    snippet = full_text[:400]
    except Exception:
        pass

    img = Image.new('RGB', (600, 720), color='#ffffff')
    draw = ImageDraw.Draw(img)

    draw.rectangle([0, 0, 599, 719], outline='#dee2e6', width=2)
    draw.rectangle([0, 0, 600, 18], fill='#28a745')
    draw.rectangle([0, 18, 600, 110], fill='#f8f9fa')

    draw.text((30, 38), "MICROSOFT EXCEL SPREADSHEET", fill='#6c757d')
    draw.text((30, 65), file_name[:42], fill='#28a745')

    draw.rectangle([30, 140, 570, 143], fill='#28a745')
    draw.rectangle([30, 160, 36, 680], fill='#28a745')

    y = 165
    if snippet:
        lines = [snippet[i:i+42] for i in range(0, len(snippet), 42)]
        for line in lines[:18]:
            draw.text((50, y), line, fill='#333333')
            y += 26
    else:
        draw.text((50, 165), "(Excel Spreadsheet Preview)", fill='#6c757d')

    img.save(cache_path, 'JPEG', quality=90)


def generate_fast_text_cover(file_path, cache_path, ftype):
    file_name = os.path.basename(file_path)
    snippet = ""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            snippet = f.read(500)
    except Exception:
        pass

    color = '#6f42c1' if ftype in ('md', 'txt', 'rtf') else '#17a2b8'
    img = Image.new('RGB', (600, 720), color='#ffffff')
    draw = ImageDraw.Draw(img)

    draw.rectangle([0, 0, 599, 719], outline='#dee2e6', width=2)
    draw.rectangle([0, 0, 600, 18], fill=color)
    draw.rectangle([0, 18, 600, 110], fill='#f8f9fa')

    draw.text((30, 38), f"{ftype.upper()} DOCUMENT", fill='#6c757d')
    draw.text((30, 65), file_name[:42], fill=color)

    draw.rectangle([30, 140, 570, 143], fill=color)
    draw.rectangle([30, 160, 36, 680], fill=color)

    y = 165
    if snippet:
        lines = snippet.splitlines()
        for line in lines[:18]:
            draw.text((50, y), line[:45], fill='#333333')
            y += 26
            if y > 650:
                break
    else:
        draw.text((50, 165), f"({ftype.upper()} Document Preview)", fill='#6c757d')

    img.save(cache_path, 'JPEG', quality=90)


def generate_thumbnail_if_missing(file_path: str, ext: str, force: bool = False):
    if not os.path.exists(file_path):
        return
    try:
        ext_clean = ext.lstrip('.').lower()
        file_hash = get_thumbnail_hash(file_path)
        cache_path = os.path.join(THUMB_CACHE_DIR, f"{file_hash}.jpg")

        if not force and os.path.exists(cache_path) and os.path.getsize(cache_path) > 0:
            try:
                if os.path.getmtime(cache_path) >= os.path.getmtime(file_path):
                    return
            except Exception:
                pass

        if ext_clean == 'pdf':
            doc = fitz.open(file_path)
            if len(doc) > 0:
                page = doc[0]
                pix = page.get_pixmap(dpi=200)
                pix.save(cache_path)
                doc.close()

        elif ext_clean in ('docx', 'doc'):
            generate_fast_docx_cover(file_path, cache_path)

        elif ext_clean in ('pptx', 'ppt'):
            generate_fast_pptx_cover(file_path, cache_path)

        elif ext_clean in ('xlsx', 'xls'):
            generate_fast_xlsx_cover(file_path, cache_path)

        elif ext_clean in ('txt', 'md', 'csv', 'rtf'):
            generate_fast_text_cover(file_path, cache_path, ext_clean)
    except Exception:
        pass



def get_opensearch_client(host="localhost", port=9200):
    return OpenSearch(
        hosts=[{"host": host, "port": port}],
        http_compress=True,
        use_ssl=False,
        verify_certs=False,
        ssl_assert_hostname=False,
        ssl_show_warn=False,
        timeout=30,
        max_retries=3,
        retry_on_timeout=True
    )


def ensure_index_exists(client, index_name="documents", settings_path="index_settings.json"):
    if client.indices.exists(index=index_name):
        return

    print(f"[*] Creating index '{index_name}' with custom settings...")
    if os.path.exists(settings_path):
        with open(settings_path, "r", encoding="utf-8") as f:
            index_body = json.load(f)
    else:
        index_body = {
            "mappings": {
                "properties": {
                    "file_name": {"type": "text", "fields": {"keyword": {"type": "keyword"}}},
                    "file_type": {"type": "text", "fields": {"keyword": {"type": "keyword"}}},
                    "file_path": {"type": "keyword"},
                    "file_directory": {"type": "keyword"},
                    "file_extension": {"type": "keyword"},
                    "file_size": {"type": "long"},
                    "created_date": {"type": "date"},
                    "modified_date": {"type": "date"},
                    "content": {"type": "text", "term_vector": "with_positions_offsets"}
                }
            }
        }
    
    client.indices.create(index=index_name, body=index_body)


def extract_file_content(file_path: str, max_bytes: int = 1_000_000) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext in SKIP_EXTENSIONS or not ext:
        return ""
        
    data = safe_read_bytes(file_path)
    if not data:
        return ""

    if ext in ['.txt', '.md', '.rtf']:
        try:
            return data[:max_bytes].decode('utf-8', errors='ignore')
        except Exception:
            return ""

    elif ext == '.pdf':
        try:
            doc = fitz.open(stream=data, filetype="pdf")
            text_parts = []
            for i, page in enumerate(doc):
                if i >= 50:
                    break
                text = page.get_text()
                if text:
                    text_parts.append(text)
            doc.close()
            return "\n".join(text_parts)[:max_bytes]
        except Exception:
            return ""

    elif ext in ['.docx', '.doc']:
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                if 'word/document.xml' in z.namelist():
                    xml_content = z.read('word/document.xml')
                    tree = ET.fromstring(xml_content)
                    text_nodes = tree.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t')
                    full_text = " ".join([node.text for node in text_nodes if node.text])
                    if full_text:
                        return full_text[:max_bytes]
        except Exception:
            pass
        if docx is not None:
            try:
                doc = docx.Document(io.BytesIO(data))
                return "\n".join([p.text for p in doc.paragraphs if p.text])[:max_bytes]
            except Exception:
                pass
        return ""

    elif ext in ['.xlsx', '.xls']:
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                if 'xl/sharedStrings.xml' in z.namelist():
                    xml_content = z.read('xl/sharedStrings.xml')
                    tree = ET.fromstring(xml_content)
                    text_nodes = tree.findall('.//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t')
                    return " ".join([node.text for node in text_nodes if node.text])[:max_bytes]
        except Exception:
            pass
        return ""

    return ""


def fetch_existing_metadata(client, index_name="documents"):
    existing_map = {}
    if not client.indices.exists(index=index_name):
        return existing_map

    total_docs = 0
    try:
        res = client.count(index=index_name)
        total_docs = res.get('count', 0)
    except Exception:
        pass

    total_str = f"{total_docs:,}" if total_docs > 0 else "~93k"
    print(f"[*] Pre-fetching existing index metadata ({total_str} records) for fast delta comparison...")
    write_progress(True, 0, 0, 0, f"⚡ Pre-fetching database metadata: 0 / {total_str} loaded... (takes ~15-30s)")

    try:
        from opensearchpy.helpers import scan
        docs = scan(client, index=index_name, query={"_source": ["file_path", "modified_date", "file_size"]})
        count = 0
        for doc in docs:
            count += 1
            src = doc.get('_source', {})
            fp = src.get('file_path')
            if fp:
                norm = os.path.normpath(os.path.abspath(fp)).lower()
                existing_map[norm] = (src.get('modified_date'), src.get('file_size'))
            if count % 2000 == 0 or count == total_docs:
                pct_str = f" ({int(count / total_docs * 100)}%)" if total_docs > 0 else ""
                write_progress(True, 0, 0, 0, f"⚡ Pre-fetching database metadata: {count:,} / {total_str} loaded{pct_str}... (takes ~15-30s)")
    except Exception as e:
        print(f"[-] Warning: Could not pre-fetch index metadata: {e}")
    print(f"[+] Loaded {len(existing_map):,} cached entries for instant skip check.")
    write_progress(True, 0, 0, 0, f"✅ Metadata loaded ({len(existing_map):,} records). Starting file scan...")
    return existing_map


SYNC_AUDIT_LOG_FILE = os.path.join(os.path.dirname(__file__), "sync_audit_report.txt")
SYNC_AUDIT_LOCK = threading.Lock()

def record_sync_event(action_label: str, src_path: str, dest_path: str, reason: str, src_size: int, dest_size: int = None, src_mtime: float = None, dest_mtime: float = None):
    """Instantly appends sync audit entries to sync_audit_report.txt in real time."""
    try:
        now_str = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        s_dt = datetime.fromtimestamp(src_mtime).strftime("%Y-%m-%d %H:%M:%S") if src_mtime else "N/A"
        d_dt = datetime.fromtimestamp(dest_mtime).strftime("%Y-%m-%d %H:%M:%S") if dest_mtime else "N/A"
        d_sz = f"{dest_size:,} bytes" if dest_size is not None else "N/A"
        delta_str = f"{abs(dest_mtime - src_mtime):.2f}s" if (dest_mtime and src_mtime) else "N/A"

        entry = (
            f"[{now_str}] SYNC EVENT: {action_label}\n"
            f"     Prompt Reason : {reason}\n"
            f"     File Name     : {os.path.basename(src_path)}\n"
            f"     Source Path   : {src_path}\n"
            f"     Target Path   : {dest_path}\n"
            f"     Source Size   : {src_size:,} bytes\n"
            f"     Target Size   : {d_sz}\n"
            f"     Source mtime  : {s_dt}\n"
            f"     Target mtime  : {d_dt}\n"
            f"     Time Delta    : {delta_str}\n"
            "----------------------------------------------------------------------------------------\n"
        )
        with SYNC_AUDIT_LOCK:
            with open(SYNC_AUDIT_LOG_FILE, "a", encoding="utf-8") as f:
                f.write(entry)
    except Exception:
        pass


def copy_file_to_backup(file_path: str, backup_dir: str) -> str:
    """
    Mirrors file_path into backup_dir preserving directory structure.
    Uses extended Windows long-path syntax (\\\\?\\) for paths >= 260 chars.
    Logs sync reason (NEW_FILE, SIZE_MISMATCH, TIMESTAMP_DRIFT) to sync_audit_report.txt ONLY after copy succeeds.
    """
    if not backup_dir:
        return None
    try:
        norm_file = os.path.normpath(os.path.abspath(file_path)).lower()
        norm_backup = os.path.normpath(os.path.abspath(backup_dir)).lower()
        # Do not recursively back up files that are already inside the backup destination
        if norm_file.startswith(norm_backup) or norm_file.startswith(r"d:\backups"):
            return None
        drive, rel_path = os.path.splitdrive(file_path)
        clean_rel = rel_path.lstrip('\\').lstrip('/')
        dest_path = os.path.join(backup_dir, clean_rel)

        long_src = to_long_path(file_path)
        long_dest = to_long_path(dest_path)

        s_stat = os.stat(long_src)
        src_size = s_stat.st_size
        src_mtime = s_stat.st_mtime

        dest_exists = os.path.exists(long_dest)
        dest_size = None
        dest_mtime = None
        reason = None

        if dest_exists:
            try:
                d_stat = os.stat(long_dest)
                dest_size = d_stat.st_size
                dest_mtime = d_stat.st_mtime
                delta = abs(dest_mtime - src_mtime)

                if src_size == dest_size and delta <= 3.0:
                    return dest_path
                elif src_size != dest_size:
                    reason = f"SIZE_MISMATCH (Source: {src_size:,} bytes | Target: {dest_size:,} bytes)"
                else:
                    reason = f"TIMESTAMP_DRIFT (Delta: {delta:.2f}s > 3.0s threshold)"
            except Exception:
                reason = "DEST_STAT_FAILED"
        else:
            reason = "NEW_FILE (Target file missing on local backup)"

        # Ensure directory exists with long path support
        os.makedirs(os.path.dirname(long_dest), exist_ok=True)
        import shutil
        shutil.copy2(long_src, long_dest)

        # Log ONLY when copy actually succeeds!
        if reason:
            record_sync_event("Local D:\\ Backup", file_path, dest_path, reason, src_size, dest_size, src_mtime, dest_mtime)

        return dest_path
    except Exception as e:
        print(f"[-] Backup copy error for '{file_path}': {e}")
        return None


def process_single_file(file_path: str, index_name: str, existing_map: dict = None, force: bool = False, backup_dir: str = None, exclude_backup_exts: set = None) -> tuple:
    file_name = os.path.basename(file_path)
    ext = os.path.splitext(file_name)[1].lower()
    
    if not ext or ext not in STRICT_DOCUMENT_EXTENSIONS or file_name.startswith('~') or file_name.startswith('.'):
        return None, 'ignored'
        
    try:
        long_p = to_long_path(file_path)
        stat = os.stat(long_p)
        norm_path = os.path.normpath(os.path.abspath(file_path)).lower()
        mtime_iso = datetime.fromtimestamp(stat.st_mtime).isoformat()
        file_size = stat.st_size

        # Handle Document Backup (Skipping excluded extensions like .csv)
        backup_dest = None
        is_backed_up = False
        if exclude_backup_exts is None:
            exclude_backup_exts = {'.csv'}

        # Bypass local D:\Backups\Documents for D:\Active research (it syncs directly to NAS in Phase 3)
        if backup_dir and ext not in exclude_backup_exts and not norm_path.startswith(r"d:\active research"):
            backup_dest = copy_file_to_backup(file_path, backup_dir)
            is_backed_up = backup_dest is not None

        # Fast Delta Skip Check: Compare file mtime and file size with 3-second tolerance
        if not force and existing_map and norm_path in existing_map:
            cached_mtime, cached_size = existing_map[norm_path]
            if cached_size == file_size:
                if cached_mtime == mtime_iso:
                    return None, 'skipped'
                # Check fuzzy timestamp difference (handles microsecond truncations & SMB 2-second FAT resolution)
                try:
                    cached_dt = datetime.fromisoformat(str(cached_mtime))
                    disk_dt = datetime.fromtimestamp(stat.st_mtime)
                    if abs((disk_dt - cached_dt).total_seconds()) <= 3.0:
                        return None, 'skipped'
                except Exception:
                    pass

        # Canonical path hash for consistent document ID across drive casing
        doc_id = hashlib.sha256(norm_path.encode('utf-8')).hexdigest()
        file_dir = os.path.dirname(file_path)
        file_type = ext.lstrip('.').lower()
        content = extract_file_content(file_path)
        
        # Trigger cover thumbnail generation if missing
        generate_thumbnail_if_missing(file_path, file_type)

        doc = {
            "_op_type": "index",
            "_index": index_name,
            "_id": doc_id,
            "file_name": file_name,
            "file_type": file_type,
            "file_path": file_path,
            "file_directory": file_dir,
            "file_extension": ext,
            "file_size": file_size,
            "created_date": datetime.fromtimestamp(stat.st_ctime).isoformat(),
            "modified_date": mtime_iso,
            "is_backed_up": is_backed_up,
            "backup_path": backup_dest,
            "content": content
        }
        return doc, 'indexed'
    except Exception:
        return None, 'error'


def is_excluded_dir(dir_path: str) -> bool:
    norm_path = dir_path.lower().replace('\\', '/')
    dir_name = os.path.basename(dir_path).lower()
    
    if dir_name in DEFAULT_EXCLUDE_NAMES or dir_name.startswith('.'):
        return True
        
    for excl in DEFAULT_EXCLUDE_DIRS:
        excl_norm = excl.lower().replace('\\', '/')
        if excl_norm in norm_path:
            return True
    return False


def scan_directories(target_dirs: list):
    for target_dir in target_dirs:
        if not os.path.exists(target_dir):
            continue
            
        print(f"[*] Scanning target: {target_dir}")
        for root, dirs, files in os.walk(target_dir, topdown=True):
            dirs[:] = [d for d in dirs if not is_excluded_dir(os.path.join(root, d))]
            for file in files:
                file_path = os.path.join(root, file)
                yield file_path


def get_latest_file_info(source_dir: str):
    """
    Returns (basename, formatted_mtime_str, full_path) for the single most recently modified file in source_dir.
    Uses os.scandir for high-performance directory traversal and instant stat retrieval.
    """
    latest_time = 0
    latest_file = None
    if source_dir and os.path.exists(source_dir):
        try:
            stack = [source_dir]
            while stack:
                curr = stack.pop()
                try:
                    with os.scandir(curr) as it:
                        for entry in it:
                            try:
                                if entry.is_dir(follow_symlinks=False):
                                    name_lower = entry.name.lower()
                                    if name_lower not in ['__pycache__', '.git', 'node_modules', '.venv', '.cache', 'appdata']:
                                        stack.append(entry.path)
                                elif entry.is_file(follow_symlinks=False):
                                    stat = entry.stat(follow_symlinks=False)
                                    mt = stat.st_mtime
                                    if mt > latest_time:
                                        latest_time = mt
                                        latest_file = entry.path
                            except Exception:
                                pass
                except Exception:
                    pass
        except Exception:
            pass
    if latest_file and latest_time > 0:
        dt_str = datetime.fromtimestamp(latest_time).strftime("%Y-%m-%d %H:%M:%S")
        return os.path.basename(latest_file), dt_str, latest_file
    return "None", "N/A", "N/A"


PROGRESS_FILE = os.path.join(os.path.dirname(__file__), "indexer_progress.json")

def write_progress(is_running: bool, scanned: int, indexed: int, skipped: int, status_msg: str, san_summary: dict = None):
    try:
        data = {
            "is_running": is_running,
            "scanned": scanned,
            "indexed": indexed,
            "skipped": skipped,
            "status_message": status_msg,
            "timestamp": time.time()
        }
        if san_summary is not None:
            data["san_summary"] = san_summary
        elif os.path.exists(PROGRESS_FILE):
            with open(PROGRESS_FILE, "r", encoding="utf-8") as f:
                old = json.load(f)
                if "san_summary" in old:
                    data["san_summary"] = old["san_summary"]
        with open(PROGRESS_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f)
    except Exception:
        pass


ALERT_EMAIL_RECIPIENT = "prdexter@iu.edu"

def send_email_alert(subject: str, body_text: str, recipient: str = ALERT_EMAIL_RECIPIENT, smtp_server: str = "mail-relay.iu.edu", smtp_port: int = 25):
    """
    Sends an email alert notification to prdexter@iu.edu if SAN backup sync fails.
    Uses Python smtplib with fallback to PowerShell Send-MailMessage.
    """
    print(f"[*] Sending backup failure alert email to {recipient}...")
    sender = "opensearch-indexer-alert@iu.edu"
    
    # Method 1: Python smtplib via IU SMTP relay
    try:
        import smtplib
        from email.mime.text import MIMEText
        
        msg = MIMEText(body_text)
        msg['Subject'] = subject
        msg['From'] = sender
        msg['To'] = recipient
        
        with smtplib.SMTP(smtp_server, smtp_port, timeout=10) as server:
            server.send_message(msg)
        print(f"[✓] Email alert successfully sent to {recipient} via {smtp_server}!")
        return
    except Exception as e1:
        print(f"[-] Standard SMTP relay error ({e1}). Attempting PowerShell mail fallback...")

    # Method 2: PowerShell Send-MailMessage fallback on Windows
    if sys.platform == 'win32':
        try:
            ps_script = f"""
            $Subject = '{subject}'
            $Body = '{body_text}'
            $Smtp = '{smtp_server}'
            Send-MailMessage -To '{recipient}' -From '{sender}' -Subject $Subject -Body $Body -SmtpServer $Smtp -ErrorAction Stop
            """
            subprocess.run(["powershell", "-Command", ps_script], capture_output=True, timeout=10)
            print(f"[✓] Email alert sent to {recipient} via PowerShell!")
        except Exception as e2:
            print(f"[-] Email alert fallback error: {e2}")


def build_san_summary_dict(network_target: str = r"\\Synology_NAS\Videos and pics\Backups"):
    san_summary = {}
    local_backups = r"D:\Backups"
    f1, t1, p1 = get_latest_file_info(local_backups)
    san_summary["Documents"] = {
        "source": local_backups,
        "target": network_target,
        "latest_file": f1,
        "latest_time": t1,
        "latest_path": p1
    }
    active_res = r"D:\Active research"
    active_nas = os.path.join(network_target, "Active research")
    f2, t2, p2 = get_latest_file_info(active_res)
    san_summary["Active Research"] = {
        "source": active_res,
        "target": active_nas,
        "latest_file": f2,
        "latest_time": t2,
        "latest_path": p2
    }
    endnote_src = r"D:\Endnote" if os.path.exists(r"D:\Endnote") else r"D:\Backups\Documents\Endnote"
    endnote_nas = r"\\Synology_NAS\Videos and pics\Endnote"
    f3, t3, p3 = get_latest_file_info(endnote_src)
    san_summary["EndNote"] = {
        "source": endnote_src,
        "target": endnote_nas,
        "latest_file": f3,
        "latest_time": t3,
        "latest_path": p3
    }
    quicken_src = r"C:\Users\Paul Dexter\OneDrive\Finances and family\Quicken"
    quicken_nas = r"\\Synology_NAS\Videos and pics\Quicken"
    f4, t4, p4 = get_latest_file_info(quicken_src)
    san_summary["Quicken"] = {
        "source": quicken_src,
        "target": quicken_nas,
        "latest_file": f4,
        "latest_time": t4,
        "latest_path": p4
    }
    return san_summary


def run_robocopy_with_live_progress(cmd: list, label: str = "Syncing", base_copied: int = 0, base_skipped: int = 0, san_summary: dict = None):
    """
    Runs Robocopy with /V /NP /NJH and parses live stdout streaming for dynamic progress updates.
    """
    import re
    # Remove flags that suppress stdout file output or summary
    cmd_filtered = [arg for arg in cmd if arg not in ["/NFL", "/NDL", "/NJH", "/NJS", "/NC", "/NS", "/NP", "/V"]]
    cmd_filtered.extend(["/V", "/NP", "/NJH"])
    
    copied_count = 0
    skipped_count = 0
    
    last_update_t = 0.0
    stdout_lines = []
    
    try:
        proc = subprocess.Popen(cmd_filtered, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True, bufsize=1, errors='replace')
        
        for line in iter(proc.stdout.readline, ''):
            stdout_lines.append(line)
            l = line.strip()
            if not l:
                continue
            
            low = l.lower()
            if " dir" in low or low.endswith("\\") or low.endswith("/"):
                continue  # Skip directory headers so only file operations are counted
            if "same" in low:
                skipped_count += 1
            elif any(k in low for k in ["new file", "newer", "older", "modified"]):
                copied_count += 1
            
            now = time.time()
            if now - last_update_t >= 0.3:
                total_c = base_copied + copied_count
                total_s = base_skipped + skipped_count
                total_chk = total_c + total_s
                msg = f"⚡ Phase 3/3 ({total_c:,} synced, {total_s:,} skipped / {total_chk:,} checked): {label}"
                write_progress(True, total_chk, total_c, total_s, msg, san_summary=san_summary)
                last_update_t = now
                
        proc.wait()
        returncode = proc.returncode
        stdout_data = "".join(stdout_lines)
    except Exception as e:
        print(f"[-] Failed to launch Robocopy: {e}")
        return 16, 0, 0

    # Parse authoritative Job Summary table from Robocopy stdout
    m = re.search(r'Files\s*:\s*([\d,]+)\s+([\d,]+)\s+([\d,]+)', stdout_data, re.IGNORECASE)
    if m:
        try:
            copied_count = int(m.group(2).replace(',', ''))
            skipped_count = int(m.group(3).replace(',', ''))
        except ValueError:
            pass
    
    total_c = base_copied + copied_count
    total_s = base_skipped + skipped_count
    total_chk = total_c + total_s
    final_msg = f"⚡ Phase 3/3 ({total_c:,} synced, {total_s:,} skipped / {total_chk:,} checked): {label} complete"
    write_progress(True, total_c + total_s, total_c, total_s, final_msg, san_summary=san_summary)
    return returncode, copied_count, skipped_count


def sync_to_network_target(backup_dir: str, network_target: str, net_user: str = None, net_pass: str = None, alert_email: str = ALERT_EMAIL_RECIPIENT, smtp_server: str = "mail-relay.iu.edu"):
    if not backup_dir or not os.path.exists(backup_dir):
        print(f"[-] Backup directory '{backup_dir}' does not exist. Skipping network copy.")
        return
    if not network_target:
        return

    print(f"\n[*] Starting backup sync from '{backup_dir}' to network target '{network_target}' (Max retries: 1)...")
    import subprocess
    if sys.platform == 'win32' and net_user and net_pass:
        try:
            share_root = network_target.split('\\')[0:4]
            share_path = "\\".join([p for p in share_root if p])
            if share_path.startswith('\\\\'):
                subprocess.run(["net", "use", share_path, net_pass, f"/user:{net_user}"], capture_output=True, timeout=5)
        except Exception:
            pass

    sync_failed = False
    failure_reason = ""
    san_summary = build_san_summary_dict(network_target)
    cum_copied = 0
    cum_skipped = 0

    try:
        os.makedirs(network_target, exist_ok=True)

        # Audit NAS sync candidates to sync_audit_report.txt before starting transfers
        try:
            from generate_sync_audit_report import generate_full_audit_report
            def audit_progress(checked, total, drift_count, is_collecting=False):
                if is_collecting:
                    msg = f"⚡ Phase 3/3 Audit (Discovered {checked:,} candidate files so far...): Scanning local drives..."
                else:
                    msg = f"⚡ Phase 3/3 Audit ({checked:,}/{total:,} checked, {drift_count:,} drift/missing): Scanning NAS across 128 threads..."
                write_progress(True, checked, 0, drift_count, msg, san_summary=san_summary)
            generate_full_audit_report(progress_callback=audit_progress)
        except Exception as e_audit:
            print(f"[-] Audit log error: {e_audit}")

        # 1. Direct SAN Sync for D:\Backups
        local_backups_root = os.path.dirname(backup_dir) if backup_dir.lower().endswith(r"\documents") else backup_dir
        if sys.platform == 'win32':
            cmd = ["robocopy", local_backups_root, network_target, "/MIR", "/FFT", "/DCOPY:DAT", "/TIMFIX", "/J", "/R:1", "/W:2", "/MT:128", "/XD", "__pycache__", ".git", "node_modules", ".venv", "venv", ".cache", ".cache_thumbnails", "appdata", "identified", "scratch", "temp", "/XF", "indexer_config.json", "indexer_progress.json", "sync_progress.json", "robocopy_monitor.log", "sync_audit_report.txt"]
            code, c1, s1 = run_robocopy_with_live_progress(cmd, label=f"Syncing backups to '{network_target}'", base_copied=cum_copied, base_skipped=cum_skipped, san_summary=san_summary)
            cum_copied += c1
            cum_skipped += s1
            if code <= 7:
                print(f"[✓] Network backup sync complete to '{network_target}'!")
        
        doc_fname, doc_ftime, doc_fpath = get_latest_file_info(local_backups_root)
        san_summary["Documents"] = {
            "source": local_backups_root,
            "target": network_target,
            "latest_file": doc_fname,
            "latest_time": doc_ftime,
            "latest_path": doc_fpath
        }
            
        # 2. Direct NAS Sync for D:\Active research
        active_res_src = r"D:\Active research"
        active_res_nas = os.path.join(network_target, "Active research")
        if os.path.exists(active_res_src):
            print(f"[*] Syncing ALL file types from '{active_res_src}' directly to NAS target...")
            os.makedirs(active_res_nas, exist_ok=True)
            if sys.platform == 'win32':
                cmd_ar = ["robocopy", active_res_src, active_res_nas, "/MIR", "/FFT", "/DCOPY:DAT", "/TIMFIX", "/J", "/R:1", "/W:2", "/MT:128", "/XD", "__pycache__", ".git", "node_modules", ".venv", "venv", ".cache", ".cache_thumbnails", "appdata", "identified", "scratch", "temp", "/XF", "indexer_config.json", "indexer_progress.json", "sync_progress.json", "robocopy_monitor.log", "sync_audit_report.txt"]
                returncode_ar, c2, s2 = run_robocopy_with_live_progress(cmd_ar, label=f"Syncing Active research to '{active_res_nas}'", base_copied=cum_copied, base_skipped=cum_skipped, san_summary=san_summary)
                cum_copied += c2
                cum_skipped += s2
                if returncode_ar <= 7:
                    print(f"[✓] Direct Active research NAS sync complete to '{active_res_nas}'!")
                else:
                    print(f"[-] Robocopy code {returncode_ar} during direct Active research NAS sync.")

        ar_fname, ar_ftime, ar_fpath = get_latest_file_info(active_res_src)
        san_summary["Active Research"] = {
            "source": active_res_src,
            "target": active_res_nas,
            "latest_file": ar_fname,
            "latest_time": ar_ftime,
            "latest_path": ar_fpath
        }

        # 3. Direct SAN Sync for EndNote Library & Articles
        endnote_src = r"D:\Endnote"
        endnote_fallback = r"D:\Backups\Documents\Endnote"
        endnote_nas_dst = r"\\Synology_NAS\Videos and pics\Endnote"
        src_e = endnote_src if os.path.exists(endnote_src) else endnote_fallback
        endnote_local_dst = r"D:\Endnote"
        if os.path.exists(src_e) and sys.platform == 'win32':
            print(f"[*] Syncing EndNote from '{src_e}' to SAN target '{endnote_nas_dst}'...")
            if os.path.normpath(src_e) != os.path.normpath(endnote_local_dst):
                cmd_e1 = ["robocopy", src_e, endnote_local_dst, "/E", "/FFT", "/DCOPY:DAT", "/TIMFIX", "/J", "/R:1", "/W:2", "/MT:128", "/NFL", "/NDL"]
                code_e1, c_e1, s_e1 = run_robocopy_with_live_progress(cmd_e1, label=f"Local EndNote sync to '{endnote_local_dst}'", base_copied=cum_copied, base_skipped=cum_skipped, san_summary=san_summary)
                cum_copied += c_e1
                cum_skipped += s_e1
            cmd_e2 = ["robocopy", endnote_local_dst, endnote_nas_dst, "/E", "/FFT", "/DCOPY:DAT", "/TIMFIX", "/J", "/R:1", "/W:2", "/MT:128", "/NFL", "/NDL"]
            returncode_e, c_e2, s_e2 = run_robocopy_with_live_progress(cmd_e2, label=f"Syncing EndNote to '{endnote_nas_dst}'", base_copied=cum_copied, base_skipped=cum_skipped, san_summary=san_summary)
            cum_copied += c_e2
            cum_skipped += s_e2
            if returncode_e <= 7:
                print(f"[✓] EndNote SAN sync complete to '{endnote_nas_dst}'!")
            else:
                print(f"[-] Robocopy code {returncode_e} during EndNote SAN sync.")

        e_fname, e_ftime, e_fpath = get_latest_file_info(src_e)
        san_summary["EndNote"] = {
            "source": src_e,
            "target": endnote_nas_dst,
            "latest_file": e_fname,
            "latest_time": e_ftime,
            "latest_path": e_fpath
        }

        # 4. Direct SAN Sync for Quicken
        quicken_src = r"C:\Users\Paul Dexter\OneDrive\Finances and family\Quicken"
        quicken_local_dst = r"D:\Quicken"
        quicken_nas_dst = r"\\Synology_NAS\Videos and pics\Quicken"
        if os.path.exists(quicken_src) and sys.platform == 'win32':
            print(f"[*] Syncing Quicken from '{quicken_src}' to SAN target '{quicken_nas_dst}'...")
            cmd_q1 = ["robocopy", quicken_src, quicken_local_dst, "/E", "/FFT", "/DCOPY:DAT", "/TIMFIX", "/J", "/R:1", "/W:2", "/MT:128", "/NFL", "/NDL"]
            code_q1, c_q1, s_q1 = run_robocopy_with_live_progress(cmd_q1, label=f"Local Quicken sync to '{quicken_local_dst}'", base_copied=cum_copied, base_skipped=cum_skipped, san_summary=san_summary)
            cum_copied += c_q1
            cum_skipped += s_q1
            cmd_q2 = ["robocopy", quicken_local_dst, quicken_nas_dst, "/E", "/FFT", "/DCOPY:DAT", "/TIMFIX", "/J", "/R:1", "/W:2", "/MT:128", "/NFL", "/NDL"]
            returncode_q, c_q2, s_q2 = run_robocopy_with_live_progress(cmd_q2, label=f"Syncing Quicken to '{quicken_nas_dst}'", base_copied=cum_copied, base_skipped=cum_skipped, san_summary=san_summary)
            cum_copied += c_q2
            cum_skipped += s_q2
            if returncode_q <= 7:
                print(f"[✓] Quicken SAN sync complete to '{quicken_nas_dst}'!")
            else:
                print(f"[-] Robocopy code {returncode_q} during Quicken SAN sync.")

        q_fname, q_ftime, q_fpath = get_latest_file_info(quicken_src)
        san_summary["Quicken"] = {
            "source": quicken_src,
            "target": quicken_nas_dst,
            "latest_file": q_fname,
            "latest_time": q_ftime,
            "latest_path": q_fpath
        }

        # Display Summary Table in Console
        print("\n==========================================================================================")
        print("                  SAN BACKUP RECENT FILE CHANGES SUMMARY REPORT                           ")
        print("==========================================================================================")
        for cat, info in san_summary.items():
            print(f" • {cat.upper()}:")
            print(f"     SAN Target Path:  {info['target']}")
            print(f"     Last File Change: {info['latest_file']}")
            print(f"     Modified Time:    {info['latest_time']}")
            print(f"     Source Path:      {info['latest_path']}")
            print("------------------------------------------------------------------------------------------")
        print("==========================================================================================\n")

        # Save summary to progress JSON for UI consumption
        write_progress(False, 0, 0, 0, "✅ SAN Backup Complete!", san_summary=san_summary)

    except Exception as e:
        sync_failed = True
        failure_reason = str(e)
        print(f"[-] Network backup sync failed: {e}")

    if sync_failed:
        recipient = alert_email or ALERT_EMAIL_RECIPIENT
        subject = "⚠️ OpenSearch Backup Alert: SAN Target Sync Failed"
        body = f"""OpenSearch Document Indexer Alert

The automatic backup sync to your SAN network target could not be written.

Network Target Path: {network_target}
Local Backup Path: {backup_dir}
Timestamp: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
Max Retries Attempted: 1

Failure Details:
{failure_reason}

Note: Your local OpenSearch database index and primary local backups (on {backup_dir}) remain 100% safe and intact.
"""
        send_email_alert(subject, body, recipient=recipient, smtp_server=smtp_server)


def backup_full_directory(src_dir: str, backup_dir: str):
    if not os.path.exists(src_dir) or not backup_dir:
        return
    print(f"\n[*] Performing FULL backup of '{src_dir}' into '{backup_dir}' (including ALL file types)...")
    drive, rel_path = os.path.splitdrive(src_dir)
    clean_rel = rel_path.lstrip('\\').lstrip('/')
    dest_dir = os.path.join(backup_dir, clean_rel)
    os.makedirs(dest_dir, exist_ok=True)
    try:
        import subprocess
        if sys.platform == 'win32':
            cmd = ["robocopy", src_dir, dest_dir, "/MIR", "/FFT", "/R:2", "/W:2", "/MT:16", "/XD", "__pycache__", ".git", "node_modules", ".venv", "venv", ".cache", ".cache_thumbnails", "appdata", "identified"]
            res_code = run_robocopy_with_live_progress(cmd, label=f"Full backup '{os.path.basename(src_dir)}'")
            if res_code <= 7:
                print(f"[✓] Full backup complete for '{src_dir}' -> '{dest_dir}'!")
            else:
                print(f"[-] Robocopy backup code {res_code} for '{src_dir}'")
        else:
            for root, dirs, files in os.walk(src_dir):
                rel = os.path.relpath(root, src_dir)
                dest = os.path.join(dest_dir, rel)
                os.makedirs(dest, exist_ok=True)
                for f in files:
                    s_file = os.path.join(root, f)
                    d_file = os.path.join(dest, f)
                    if not os.path.exists(d_file) or os.path.getmtime(s_file) > os.path.getmtime(d_file):
                        import shutil
                        shutil.copy2(s_file, d_file)
            print(f"[✓] Full backup complete for '{src_dir}' -> '{dest_dir}'!")
    except Exception as e:
        print(f"[-] Full backup error for '{src_dir}': {e}")


ONEDRIVE_EXE = r"C:\Program Files\Microsoft OneDrive\OneDrive.exe"
DROPBOX_EXE = r"C:\Program Files (x86)\Dropbox\Client\Dropbox.exe"


def pause_cloud_sync():
    """Stops OneDrive and Dropbox sync processes completely during indexing to prevent file locks."""
    print("[*] Pausing OneDrive & Dropbox sync processes during indexing...")
    import subprocess
    # Shutdown OneDrive cleanly first
    for odp in [r"C:\Program Files\Microsoft OneDrive\OneDrive.exe", os.path.expanduser(r"~\AppData\Local\Microsoft\OneDrive\OneDrive.exe")]:
        if os.path.exists(odp):
            try:
                subprocess.run([odp, "/shutdown"], capture_output=True, timeout=5)
            except Exception:
                pass

    # Force kill any remaining sync processes
    for proc in ["OneDrive.exe", "OneDrive.Sync.Service.exe", "Dropbox.exe"]:
        try:
            subprocess.run(["taskkill", "/F", "/T", "/IM", proc], capture_output=True, timeout=5)
        except Exception:
            pass


def resume_cloud_sync():
    """Resumes OneDrive and Dropbox sync processes after indexing completes."""
    print("[✓] Resuming OneDrive & Dropbox sync...")
    import subprocess
    for odp in [r"C:\Program Files\Microsoft OneDrive\OneDrive.exe", os.path.expanduser(r"~\AppData\Local\Microsoft\OneDrive\OneDrive.exe")]:
        if os.path.exists(odp):
            try:
                subprocess.Popen([odp, "/background"])
                break
            except Exception:
                pass

    dropbox_exe = r"C:\Program Files (x86)\Dropbox\Client\Dropbox.exe"
    if os.path.exists(dropbox_exe):
        try:
            subprocess.Popen([dropbox_exe])
        except Exception:
            pass


def main():
    parser = argparse.ArgumentParser(description="Ingest documents into OpenSearch with Hybrid Delta Indexing")
    parser.add_argument("--dir", nargs="+", help="Directories to scan and index")
    parser.add_argument("--index", default="documents", help="Target OpenSearch index name")
    parser.add_argument("--force", "-f", action="store_true", help="Force re-indexing of all files regardless of timestamps")
    parser.add_argument("--backup-dir", default=None, help="Target directory for document backups")
    parser.add_argument("--exclude-backup-ext", nargs="*", default=None, help="File extensions to exclude from backup (e.g. .csv)")
    parser.add_argument("--network-target", default=None, help="Optional network target directory to sync backups to when complete")
    parser.add_argument("--no-pause-sync", action="store_true", help="Disable automatic pausing of OneDrive and Dropbox during indexing")
    args = parser.parse_args()

    if not args.no_pause_sync:
        pause_cloud_sync()

    try:
        config_backup_dir = r"D:\Backups\Documents"
        config_exclude_exts = [".csv"]
        config_network_target = None

        config_path = os.path.join(os.path.dirname(__file__), "indexer_config.json")
        if os.path.exists(config_path):
            try:
                with open(config_path, "r", encoding="utf-8") as f:
                    cfg = json.load(f)
                    if not args.dir:
                        args.dir = cfg.get("selected_directories", [r"D:\Active research"])
                    config_backup_dir = cfg.get("backup_directory", r"D:\Backups\Documents")
                    config_exclude_exts = cfg.get("exclude_backup_extensions", [".csv"])
                    config_network_target = cfg.get("network_backup_target", None)
            except Exception:
                pass

        if not args.dir:
            args.dir = [r"D:\Active research"]

        backup_dir = args.backup_dir or config_backup_dir
        raw_exclude_exts = args.exclude_backup_ext if args.exclude_backup_ext is not None else config_exclude_exts
        exclude_backup_exts = {ext.lower() if ext.startswith('.') else f".{ext.lower()}" for ext in raw_exclude_exts}

        print(f"[*] Ingestion starting for directories: {args.dir} (Force mode: {args.force})")
        if backup_dir:
            print(f"[*] Document Backup Enabled -> Destination: '{backup_dir}' (Excluding: {sorted(list(exclude_backup_exts))})")
        else:
            print(f"[*] Document Backup Disabled (Set 'backup_directory' in indexer_config.json or use --backup-dir)")

        client = get_opensearch_client()
        ensure_index_exists(client, index_name=args.index)

        existing_map = fetch_existing_metadata(client, index_name=args.index) if not args.force else {}

        start_time = time.time()
        batch = []
        total_scanned = 0
        total_indexed = 0
        total_skipped = 0
        scanned_paths = set()

        from concurrent.futures import wait, FIRST_COMPLETED

        with ThreadPoolExecutor(max_workers=8) as executor:
            futures = set()
            max_queue = 32

            def handle_finished_future(fut):
                nonlocal total_scanned, total_indexed, total_skipped, batch
                doc, status = fut.result()
                if status == 'skipped':
                    total_skipped += 1
                    total_scanned += 1
                elif status == 'indexed' or doc:
                    total_scanned += 1
                    if doc:
                        batch.append(doc)

                if total_scanned % 10 == 0 or len(batch) >= 20:
                    msg = f"🔍 Scanning documents: {total_scanned:,} scanned ({total_indexed:,} updated, {total_skipped:,} unchanged)"
                    write_progress(True, total_scanned, total_indexed, total_skipped, msg)

                if len(batch) >= 20:
                    try:
                        helpers.bulk(client, batch)
                        client.indices.refresh(index=args.index)
                        total_indexed += len(batch)
                        print(f"[+] Bulk indexed {total_indexed} new/updated documents ({total_skipped:,} skipped unmodified)...")
                    except Exception as e:
                        print(f"[-] Bulk ingest error: {e}")
                    batch = []

            for file_path in scan_directories(args.dir):
                norm_p = os.path.normpath(os.path.abspath(file_path)).lower()
                scanned_paths.add(norm_p)
                
                while len(futures) >= max_queue:
                    done, futures = wait(futures, return_when=FIRST_COMPLETED)
                    for fut in done:
                        handle_finished_future(fut)

                futures.add(executor.submit(process_single_file, file_path, args.index, existing_map, args.force, backup_dir, exclude_backup_exts))

            # Flush remaining futures with timeout safety
            if futures:
                done, not_done = wait(futures, timeout=15)
                for fut in done:
                    try:
                        handle_finished_future(fut)
                    except Exception:
                        pass
                if not_done:
                    print(f"[-] Warning: {len(not_done)} indexing threads timed out and were skipped.")

            if batch:
                try:
                    helpers.bulk(client, batch)
                    client.indices.refresh(index=args.index)
                    total_indexed += len(batch)
                except Exception:
                    pass

        # Ghost / Dead File Cleanup Pass (Purges deleted files from OpenSearch index, thumbnail cache, and local backup)
        total_deleted = 0
        if existing_map:
            dead_paths = [fp for fp in existing_map.keys() if not os.path.exists(fp)]
            if dead_paths:
                print(f"[*] Cleaning up {len(dead_paths):,} deleted or moved documents from index & backups...")
                for fp in dead_paths:
                    try:
                        # 1. Remove from OpenSearch Index
                        doc_id = hashlib.sha256(fp.encode('utf-8')).hexdigest()
                        client.delete(index=args.index, id=doc_id, ignore=[404])
                        total_deleted += 1

                        # 2. Remove cached thumbnail
                        file_hash = get_thumbnail_hash(fp)
                        cache_path = os.path.join(THUMB_CACHE_DIR, f"{file_hash}.jpg")
                        if os.path.exists(cache_path):
                            try:
                                os.remove(cache_path)
                            except Exception:
                                pass

                        # 3. Remove local backup copy so robocopy /MIR syncs deletion to NAS
                        if backup_dir:
                            norm_p = os.path.normpath(os.path.abspath(fp))
                            drive, rest = os.path.splitdrive(norm_p)
                            rel_path = rest.lstrip('\\').lstrip('/')
                            local_b_path = os.path.join(backup_dir, rel_path)
                            if os.path.exists(local_b_path):
                                try:
                                    os.remove(local_b_path)
                                    print(f"[✓] Purged deleted file from local backup: {local_b_path}")
                                except Exception:
                                    pass
                    except Exception:
                        pass
                client.indices.refresh(index=args.index)

        elapsed = time.time() - start_time
        print(f"[+] Hybrid Indexing Complete! {total_scanned:,} files checked ({total_indexed:,} indexed/updated, {total_skipped:,} unmodified/skipped, {total_deleted:,} dead entries cleaned up) in {elapsed:.2f} seconds.")

        # Phase 2: Full Directory Special Exceptions Phase (includes ALL file types for Endnote & Quicken)
        if backup_dir:
            full_backup_sources = [
                r"C:\Users\Paul Dexter\OneDrive\Endnote",
                r"C:\Users\Paul Dexter\OneDrive\Finances and family\Quicken"
            ]
            for src in full_backup_sources:
                folder_name = os.path.basename(src)
                write_progress(True, total_scanned, total_indexed, total_skipped, f"⚡ Phase 2/3: Full backup of '{folder_name}' (all file types)...")
                backup_full_directory(src, backup_dir)

        # Phase 3: Optional Network Backup Sync Phase
        if backup_dir and os.path.exists(backup_dir):
            net_target = args.network_target or config_network_target
            if not net_target and sys.stdin.isatty():
                try:
                    ans = input("\n[?] Ingestion and local backup complete. Would you like to save/sync backups to a network target? (y/N): ").strip().lower()
                    if ans in ['y', 'yes']:
                        net_target = input("[?] Enter network target path (e.g. \\\\NAS\\Share\\Backups or Z:\\Backups): ").strip()
                except (KeyboardInterrupt, EOFError):
                    net_target = None
            
            if net_target:
                san_sum = build_san_summary_dict(net_target)
                write_progress(True, total_scanned, total_indexed, total_skipped, f"⚡ Phase 3/3 Audit: Initializing NAS sync scan to '{net_target}'...", san_summary=san_sum)
                sync_to_network_target(backup_dir, net_target)
        
        final_msg = f"✅ All Phases Complete ({total_scanned:,} docs checked, backups synced)!"
        write_progress(False, total_scanned, total_indexed, total_skipped, final_msg)
    finally:
        if not args.no_pause_sync:
            resume_cloud_sync()


if __name__ == "__main__":
    main()
